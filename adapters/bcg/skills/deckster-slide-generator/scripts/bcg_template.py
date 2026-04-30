"""
BCG Template Workflow Helper

Creates BCG presentations by populating the official BCG master template.
Structural slides (title, end, disclaimer, dividers) come pixel-perfect from
the template layouts. Content slides use the Title Only layout with custom
shapes added via XML.

Usage:
    from bcg_template import BCGDeck

    deck = BCGDeck()
    deck.add_title_slide('My Title', 'Subtitle text', '3 March 2026')
    deck.add_content_slide('Action title as complete sentence')
    deck.add_section_divider('Section Name')
    deck.add_disclaimer()
    deck.add_end_slide()
    deck.save('20260303_My_Deck.pptx')
"""

import base64, hashlib, json, os, re, shutil, sys, zlib
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx_utils import unpack, pack, clean, add_slide, make_work_dir, find_skill_root
from layout_semantics import SEMANTIC_TO_RUNTIME_LAYOUT, semantic_behavior
from theme import normalize_theme_config

# Register XML namespaces
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}
for prefix, uri in NS.items():
    ET.register_namespace(prefix, uri)

# ============================================================
# CONSTANTS
# ============================================================

FONT = 'Trebuchet MS'

COLORS = {
    'BCG_GREEN': '29BA74',
    'DARK_GREEN': '197A56',
    'DEEP_GREEN': '03522D',
    'DARK_TEXT': '575757',
    'LIGHT_BG': 'F2F2F2',
    'WHITE': 'FFFFFF',
    'TEAL': '3EAD92',
    'LIME': 'D4DF33',
    'MED_GRAY': '6E6F73',
    'NAVY': '295E7E',
    'NEGATIVE': 'D64454',
    'PAGE_NUM': 'B0B0B0',
    'STAMP_RED': 'E71C57',
    'LIGHT_TEXT': 'F2F2F2',
    'SLIDE_BG': 'F2F2F2',   # Default slide background (gray for standard BCG)
}

# Layout background data -- populated by apply_theme() from template.json
_LAYOUT_BACKGROUNDS = {}

# Active theme config -- stored by apply_theme() for content_bounds() lookups
_ACTIVE_THEME_CONFIG = None

# Active slide context -- updated by add_content_slide()/structural slide helpers
_CURRENT_SLIDE_CONTEXT = {"deck": None, "slide": None}

# Light/dark text palettes -- auto-populated by apply_theme(), BCG defaults here
LIGHT_PALETTE = {
    'header': COLORS['BCG_GREEN'],   # Brand color for headers on light backgrounds
    'text': COLORS['DARK_TEXT'],     # Body text on light backgrounds
    'secondary': COLORS['MED_GRAY'], # Muted text on light backgrounds
}
DARK_PALETTE = {
    'header': COLORS['WHITE'],   # White for headers on dark backgrounds
    'text': COLORS['WHITE'],     # White text on dark backgrounds
    'secondary': COLORS['LIGHT_BG'], # Off-white on dark backgrounds
}


# ============================================================
# CONTRAST HELPERS (WCAG 2.0)
# ============================================================

def _hex_to_rgb(hex_color):
    """Convert hex color string to (r, g, b) tuple."""
    if not hex_color:
        return None
    h = hex_color.lstrip('#')
    if len(h) != 6:
        return None
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None

def _relative_luminance(rgb):
    """Calculate relative luminance per WCAG 2.0."""
    def _ch(c):
        v = c / 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    return 0.2126 * _ch(rgb[0]) + 0.7152 * _ch(rgb[1]) + 0.0722 * _ch(rgb[2])

def _smart_accent():
    """Return the best accent color for visible shapes/fills.

    Uses BCG_GREEN (remapped to client accent1 when a theme is active).
    If the color is too light for visibility, falls back to DARK_GREEN.
    """
    hc = COLORS['BCG_GREEN']
    rgb = _hex_to_rgb(hc)
    if rgb and _relative_luminance(rgb) > 0.6:
        return COLORS.get('DARK_GREEN', hc)
    return hc


def layout_bg(layout_key=None):
    """Return the background color for a layout (from template extraction).

    Usage:
        bg = layout_bg('title_only')  # Returns the slide background hex
        deck.add_textbox(slide, text, x, y, w, h, color=text_on(bg))
    """
    if layout_key:
        slug = layout_key.lower().replace('_', '-')
        info = _LAYOUT_BACKGROUNDS.get(slug, {})
        accent = info.get('accent_fill')
        if accent:
            return accent
        slot_bg = info.get('slot_bg')
        if slot_bg:
            return slot_bg
    return COLORS.get('SLIDE_BG', 'F2F2F2')


# Known-good content zone boundaries for BCG default split layouts.
# These values are measured from the BCG master template and match the
# pptx-service layout-content-slots.json manifest. They produce correct
# visual output and were validated across hundreds of benchmark slides.
# For ee4p templates, the system falls back to fly_zone globals.
_BCG_SPLIT_ZONES = {
    'green-highlight':    {'left_x': 0.69, 'left_w': 6.80, 'right_x': 8.20, 'right_w': 4.30},
    'd-green-highlight':  {'left_x': 0.69, 'left_w': 6.80, 'right_x': 8.20, 'right_w': 4.30},
    'green-one-third':    {'left_x': 0.69, 'left_w': 3.50, 'right_x': 4.80, 'right_w': 8.00},
    'd-green-one-third':  {'left_x': 0.69, 'left_w': 3.50, 'right_x': 4.80, 'right_w': 8.00},
    'green-left-arrow':   {'left_x': 0.69, 'left_w': 3.50, 'right_x': 5.20, 'right_w': 7.50},
    'd-green-left-arrow': {'left_x': 0.69, 'left_w': 3.50, 'right_x': 5.20, 'right_w': 7.50},
    'arrow-half':         {'left_x': 0.69, 'left_w': 5.20, 'right_x': 7.00, 'right_w': 5.40},
    'd-arrow-half':       {'left_x': 0.69, 'left_w': 5.20, 'right_x': 7.00, 'right_w': 5.40},
    'green-half':         {'left_x': 0.69, 'left_w': 5.00, 'right_x': 6.67, 'right_w': 6.67},
    'd-green-half':       {'left_x': 0.69, 'left_w': 5.00, 'right_x': 6.67, 'right_w': 6.67},
    'green-two-third':    {'left_x': 0.69, 'left_w': 7.50, 'right_x': 8.55, 'right_w': 4.78},
    'd-green-two-third':  {'left_x': 0.69, 'left_w': 7.50, 'right_x': 8.55, 'right_w': 4.78},
    'left-arrow':         {'left_x': 0.69, 'left_w': 3.50, 'right_x': 5.20, 'right_w': 7.50},
    'd-left-arrow':       {'left_x': 0.69, 'left_w': 3.50, 'right_x': 5.20, 'right_w': 7.50},
    'green-arrow-half':   {'left_x': 0.69, 'left_w': 5.20, 'right_x': 7.00, 'right_w': 5.40},
    'd-green-arrow-half': {'left_x': 0.69, 'left_w': 5.20, 'right_x': 7.00, 'right_w': 5.40},
    'd-four-column-green':{'left_x': 0.69, 'left_w': 6.80, 'right_x': 8.20, 'right_w': 4.30},
}

_PRIMARY_CONTENT_ZONE = {
    # Analysis/evidence belongs on the white/main side.
    'green-highlight': 'left',
    'd-green-highlight': 'left',
    'd-four-column-green': 'left',
    # Title-led layouts reserve the accent/title panel; authored content belongs on the right.
    'green-one-third': 'right',
    'd-green-one-third': 'right',
    'green-left-arrow': 'right',
    'd-green-left-arrow': 'right',
    'left-arrow': 'right',
    'd-left-arrow': 'right',
    'green-half': 'right',
    'd-green-half': 'right',
    'green-two-third': 'right',
    'd-green-two-third': 'right',
}


def _apply_primary_zone(result, slug):
    """Normalize x/w to the layout's primary authored-content zone.

    Split layouts still expose `left_x/left_w/right_x/right_w` so variants can
    address both sides explicitly. The primary zone remaps `x/w/x_end` so
    generic content patterns do not default into reserved title panels.
    """
    zone = _PRIMARY_CONTENT_ZONE.get(slug)
    if not zone or f"{zone}_x" not in result or f"{zone}_w" not in result:
        return result

    result['full_x'] = result['x']
    result['full_y'] = result['y']
    result['full_w'] = result['w']
    result['full_h'] = result['h']
    result['full_x_end'] = result['x_end']
    result['full_y_end'] = result['y_end']

    result['x'] = result[f'{zone}_x']
    result['w'] = result[f'{zone}_w']
    result['x_end'] = round(result['x'] + result['w'], 2)
    return result


def content_bounds(layout_key=None):
    """Return the usable content zone bounds for a layout in inches.

    For split layouts, returns the left and right zone boundaries that
    match the layout's visual accent panel. These are measured values
    that produce correct visual output.

    Args:
        layout_key: Layout name (e.g. 'green_highlight', 'title_only').
                    If None, returns global fly_zone bounds.

    Returns:
        dict with keys: x, y, w, h, x_end, y_end (all in inches).
        For split layouts, also includes:
            full_x, full_y, full_w, full_h -- full fly-zone bounds
            left_x, left_w   -- left content zone (white/content side)
            right_x, right_w -- right content zone (accent/insight side)

        On title-left split layouts (`green_one_third`, `green_left_arrow`,
        `green_half`, `green_two_third`, and detail variants), `x/w/x_end`
        are normalized to the right-side authored-content zone so generic
        patterns do not land inside the reserved title panel.

    Usage:
        from bcg_template import content_bounds
        bounds = content_bounds('green_highlight')
        left_w = bounds['left_w']    # 6.80" — max width for left-side content
        right_x = bounds['right_x']  # 8.20" — where right panel starts
    """
    # Default: fly_zone globals
    default = {
        'x': CONTENT_LEFT,
        'y': CONTENT_START_Y,
        'w': CONTENT_WIDTH,
        'h': round(CONTENT_END_Y - CONTENT_START_Y, 2),
        'x_end': round(CONTENT_LEFT + CONTENT_WIDTH, 2),
        'y_end': CONTENT_END_Y,
    }

    if layout_key is None:
        current_deck = _CURRENT_SLIDE_CONTEXT.get("deck")
        current_slide = _CURRENT_SLIDE_CONTEXT.get("slide")
        if current_deck is not None and current_slide is not None:
            try:
                return dict(current_deck.current_slide_contract(current_slide)["bounds"])
            except Exception:
                pass
        return default

    slug = layout_key.lower().replace('_', '-')

    # Priority 1: split_zones from template.json (computed during ingestion)
    if _ACTIVE_THEME_CONFIG:
        layout_geom = _ACTIVE_THEME_CONFIG.get('layout_geometry', {})
        layout_info = layout_geom.get(slug, {})
        # ee4p templates store geometry under the original layout name slug,
        # not the semantic role slug.  Resolve via semantic_layout_map.
        if not layout_info:
            sem_map = _ACTIVE_THEME_CONFIG.get('semantic_layout_map', {})
            sem_key = slug.replace('-', '_')
            original_name = sem_map.get(sem_key, '')
            if original_name:
                orig_slug = re.sub(r'[^a-z0-9]+', '-', original_name.lower()).strip('-')
                layout_info = layout_geom.get(orig_slug, {})
        cz = layout_info.get('content_zone') or {}
        sz = cz.get('split_zones')
        if sz and 'left_x' in sz and 'right_x' in sz:
            result = dict(default)
            # Clamp split zones to the fly zone so content doesn't
            # start at x=0 on ee4p templates where accent panels
            # extend to the slide edge.
            fly_left = default['x']
            fly_right = default['x_end']
            lx = max(sz['left_x'], fly_left)
            lw = sz['left_w'] - (lx - sz['left_x'])
            rx = max(sz['right_x'], fly_left)
            rw = min(sz['right_w'], fly_right - rx)
            result.update({
                'left_x': round(lx, 2),
                'left_w': round(max(lw, 0.5), 2),
                'right_x': round(rx, 2),
                'right_w': round(max(rw, 0.5), 2),
            })
            # Ensure y clears the title placeholder to prevent content
            # overlapping the title on split layouts (e.g. D. Green one third).
            # Only applies when the title is WIDE (spans across the right panel).
            # Narrow titles on the left accent panel shouldn't affect right-panel Y.
            phs = layout_info.get('placeholders', [])
            for ph in phs:
                if ph.get('type') in ('title', 'ctrTitle'):
                    tb = ph['bounds']
                    if tb.get('w_pct', 0) > 55:
                        # Wide title spans into the right panel — content must start below it
                        title_bot = (tb['y_pct'] + tb['h_pct']) / 100.0 * 7.5
                        min_y = round(title_bot + 0.10, 2)
                        if min_y > result['y']:
                            result['y'] = min_y
                            result['h'] = round(result['y_end'] - min_y, 2)
                    break
            return _apply_primary_zone(result, slug)

    # Priority 2: _BCG_SPLIT_ZONES hardcoded table.
    # Used for BCG default OR as fallback for ee4p templates whose
    # layout geometry didn't derive split zones (heuristic failure).
    # When used for a non-BCG template, scale the reference zones
    # proportionally to the template's actual fly zone.
    if slug in _BCG_SPLIT_ZONES:
        result = dict(default)
        bcg_zones = dict(_BCG_SPLIT_ZONES[slug])
        if _ACTIVE_THEME_CONFIG and _ACTIVE_THEME_CONFIG.get("source") != "built-in":
            # Scale BCG reference zones to the active template's fly zone.
            # BCG default: content_left=0.69, content_width=11.96
            bcg_left = 0.69
            bcg_width = 11.96
            tmpl_left = default['x']
            tmpl_width = default['w']
            if tmpl_width > 0 and bcg_width > 0:
                scale = tmpl_width / bcg_width
                offset = tmpl_left - bcg_left * scale
                for k in ('left_x', 'right_x'):
                    if k in bcg_zones:
                        bcg_zones[k] = round(bcg_zones[k] * scale + offset, 2)
                for k in ('left_w', 'right_w'):
                    if k in bcg_zones:
                        bcg_zones[k] = round(bcg_zones[k] * scale, 2)
        result.update(bcg_zones)
        return _apply_primary_zone(result, slug)

    # Priority 3: Global fly_zone (for full-page layouts)
    return _normalize_bounds_for_behavior(default, layout_key)


def _normalize_bounds_for_behavior(bounds, layout_key):
    """Apply layout-aware top offsets to generic bounds."""
    result = dict(bounds)
    behavior = semantic_behavior(layout_key, theme_config=_ACTIVE_THEME_CONFIG)
    if behavior["structural"] or behavior["statement"]:
        return result

    if behavior["layout_class"] == "full_page" and not behavior["detail"]:
        min_y = CONTENT_START_Y
    else:
        min_y = DETAIL_CONTENT_START_Y

    if result["y"] < min_y:
        result["y"] = min_y
        result["h"] = round(result["y_end"] - min_y, 2)
        if "full_y" in result:
            result["full_y"] = min_y
            result["full_h"] = round(result["full_y_end"] - min_y, 2)
    return result


def _content_surface_bg(layout_key):
    """Return the authored content surface background for a layout."""
    slug = str(layout_key).lower().replace("_", "-")
    info = _LAYOUT_BACKGROUNDS.get(slug, {})
    return (
        info.get("content_bg")
        or info.get("layout_bg")
        or info.get("slot_bg")
        or COLORS.get("SLIDE_BG", "F2F2F2")
    )


def _source_spec_for_bounds(bounds, behavior):
    """Return a layout-aware source placement spec or None."""
    if behavior.get("source_mode") == "none":
        return None
    source_x = bounds.get("full_x", bounds.get("x", CONTENT_LEFT))
    source_w = min(bounds.get("full_w", bounds.get("w", CONTENT_WIDTH)), 9.88)
    return {
        "x": round(source_x, 2),
        "y": SOURCE_Y,
        "w": round(source_w, 2),
        "h": 0.30,
    }


def slide_contract(slide=None, deck=None):
    """Return the active slide contract for the current deck context."""
    active_deck = deck or _CURRENT_SLIDE_CONTEXT.get("deck")
    active_slide = slide or _CURRENT_SLIDE_CONTEXT.get("slide")
    if active_deck is None or active_slide is None:
        raise RuntimeError("No active slide context. Call add_content_slide() first.")
    return active_deck.current_slide_contract(active_slide)


def text_on(bg_color):
    """Return the best text color for content on the given background.

    Usage:
        deck.add_rectangle(slide, x, y, w, h, COLORS['NAVY'])
        deck.add_textbox(slide, 'Title', x, y, w, h, color=text_on(COLORS['NAVY']))

    This is the primary way to ensure readable text on any colored shape.
    """
    return _auto_text_color(bg_color, 'text')


def card_fill_for_slide():
    """Return the best card/container fill color for the current slide background.

    On light slide backgrounds (white/gray detail slides): returns LIGHT_BG (F2F2F2)
    On dark slide backgrounds (dark-themed ee4p templates): returns WHITE (FFFFFF)

    Usage in pattern variants:
        from bcg_template import card_fill_for_slide, text_on
        card_fill = card_fill_for_slide()
        card_text = text_on(card_fill)
        deck.add_rounded_rectangle(slide, x, y, w, h, card_fill)
        deck.add_textbox(slide, text, x, y, w, h, color=card_text)
    """
    slide_bg = COLORS.get('SLIDE_BG', 'F2F2F2')
    bg_rgb = _hex_to_rgb(slide_bg)
    if bg_rgb:
        bg_lum = _relative_luminance(bg_rgb)
        if bg_lum > 0.85:
            return COLORS.get('LIGHT_BG', 'F2F2F2')
        else:
            return COLORS.get('WHITE', 'FFFFFF')
    return COLORS.get('LIGHT_BG', 'F2F2F2')


def dark_fills(n=4):
    """Return n distinct dark fill colors from the current theme.

    All returned colors have white-text contrast >= 3.5:1, suitable
    for use as shape fills with white or light text. Adapts to the
    active theme (ee4p or BCG default).

    Usage:
        fills = dark_fills(4)
        for i, fill in enumerate(fills):
            deck.add_oval(slide, x, y, w, h, fill)
            deck.add_textbox(slide, label, x, y, w, h, color=text_on(fill))
    """
    # Gather candidates from semantic theme keys (order = visual priority)
    candidates = [
        COLORS.get('CONTRAST_BG'),   # dk2 — typically darkest brand color
        COLORS.get('ACCENT1'),       # primary accent
        COLORS.get('ACCENT4'),       # secondary dark accent
        COLORS.get('ACCENT6'),       # tertiary accent
        COLORS.get('ON_SURFACE'),    # dk1 — text-dark color
        COLORS.get('ACCENT2'),       # often warm/contrasting accent
        COLORS.get('ACCENT3'),       # mid-tone accent
        COLORS.get('DEEP_GREEN'),    # BCG fallback
        COLORS.get('DARK_GREEN'),    # BCG fallback
        COLORS.get('NAVY'),          # BCG fallback
    ]

    # Deduplicate, filter to dark enough for white text
    seen = set()
    result = []
    for c in candidates:
        if not c or c in seen:
            continue
        seen.add(c)
        rgb = _hex_to_rgb(c)
        if rgb and _contrast_ratio(c, 'FFFFFF') >= 3.5:
            result.append(c)
        if len(result) >= n:
            break

    # If theme is too limited, cycle what we have
    if not result:
        result = [COLORS.get('DEEP_GREEN', '03522D')]
    base_len = len(result)
    while len(result) < n:
        result.append(result[len(result) % base_len])

    return result[:n]


def _contrast_ratio(bg_hex, fg_hex):
    """WCAG 2.0 contrast ratio between two hex colors (1.0 to 21.0)."""
    bg = _hex_to_rgb(bg_hex)
    fg = _hex_to_rgb(fg_hex)
    if bg is None or fg is None:
        return 0.0
    bl = _relative_luminance(bg)
    fl = _relative_luminance(fg)
    lighter = max(bl, fl)
    darker = min(bl, fl)
    return (lighter + 0.05) / (darker + 0.05)

def _auto_text_color(bg_color, role='text'):
    """Pick the highest-contrast text color for a given background.

    Computes WCAG contrast ratios against all candidates from both
    LIGHT_PALETTE and DARK_PALETTE, then returns the one with the
    best contrast. This handles edge cases where simple luminance
    thresholds fail (e.g., medium blues, pastel accents).

    Args:
        bg_color: Hex color of the background (e.g. '005A9C').
                  If None, uses COLORS['SLIDE_BG'] (the template's slide background).
        role: 'text' (body), 'header' (titles/accents), 'secondary' (muted).
    """
    if not bg_color:
        # Use the actual slide background color if known, not hardcoded F2F2F2
        bg_color = COLORS.get('SLIDE_BG')
        return LIGHT_PALETTE.get(role, COLORS['DARK_TEXT'])

    bg_rgb = _hex_to_rgb(bg_color)
    if bg_rgb is None:
        return LIGHT_PALETTE.get(role, COLORS['DARK_TEXT'])

    # Two-tier candidate system:
    # Tier 1: Theme-designated colors (DARK_TEXT, WHITE, palette roles).
    #         These match the template's design language. Use if contrast >= 3.0:1.
    # Tier 2: Literal extremes (000000, FFFFFF). Only used when tier 1 fails,
    #         e.g. dark templates where DARK_TEXT is remapped to FFFFFF.
    tier1 = set()
    light_val = LIGHT_PALETTE.get(role)
    dark_val = DARK_PALETTE.get(role)
    if light_val:
        tier1.add(light_val)
    if dark_val:
        tier1.add(dark_val)
    tier1.add(COLORS.get('WHITE', 'FFFFFF'))
    tier1.add(COLORS.get('DARK_TEXT', '575757'))
    tier1.add(COLORS.get('ON_SURFACE', COLORS.get('DARK_TEXT', '575757')))
    tier1.discard(None)

    tier2 = {'FFFFFF', '000000'}

    # For CHROMATIC backgrounds (not gray/neutral), prefer white.
    # On saturated colors (green, blue, red), white is ALWAYS more readable
    # than gray/black — this matches BCG's own design language.
    # BCG_GREEN (29BA74) has only 2.51:1 white contrast, but gray looks worse.
    bg_r, bg_g, bg_b = bg_rgb
    is_chromatic = not (abs(bg_r - bg_g) < 30 and abs(bg_g - bg_b) < 30)
    bg_lum = _relative_luminance(bg_rgb)

    if is_chromatic and bg_lum < 0.45:
        # Chromatic bg with lum < 0.45: always use white.
        # At lum=0.45, white contrast >= 2.1:1 — acceptable for brand colors.
        return COLORS.get('WHITE', 'FFFFFF')

    # Try tier 1 first — pick best contrast from theme-designated colors
    best_t1 = COLORS.get('DARK_TEXT', '575757')
    best_t1_ratio = 0.0
    for c in tier1:
        ratio = _contrast_ratio(bg_color, c)
        if ratio > best_t1_ratio:
            best_t1_ratio = ratio
            best_t1 = c

    # If tier 1 has adequate contrast (>= 3.0:1), use it
    if best_t1_ratio >= 3.0:
        return best_t1

    # Tier 1 failed — fall back to literal extremes (dark template scenario)
    best_color = best_t1
    best_ratio = best_t1_ratio
    for c in tier2:
        ratio = _contrast_ratio(bg_color, c)
        if ratio > best_ratio:
            best_ratio = ratio
            best_color = c

    return best_color


# Content area positioning (from BCG master Layout Guide)
# Title can wrap to 2 lines ending at ~y=1.50". Content must start below that.
CONTENT_START_Y = 2.27   # Default for content below 2-line action title (standard layouts)
CONTENT_START_Y_TIGHT = 1.65  # Use when content overflows at 2.27" — tighter fit under title
DETAIL_CONTENT_START_Y = 1.70  # Detail layouts have smaller titles (h=0.36" vs 0.51")
CONTENT_END_Y = 6.50     # Max y+h for content. Source line starts at 6.74"; leave 0.24" gap.
SOURCE_Y = 6.74          # Where the source/footnote text sits
CONTENT_LEFT = 0.69
CONTENT_WIDTH = 11.96

# BCG Master Template layout mapping (layout name -> XML filename)
LAYOUTS = {
    'title':              'slideLayout1.xml',
    'title_only':         'slideLayout2.xml',
    'title_and_text':     'slideLayout3.xml',
    'gray_slice':         'slideLayout4.xml',
    'section_header_box': 'slideLayout5.xml',
    'section_header_line':'slideLayout6.xml',
    'white_one_third':    'slideLayout7.xml',
    'green_highlight':    'slideLayout8.xml',
    'green_one_third':    'slideLayout9.xml',
    'green_half':         'slideLayout10.xml',
    'green_two_third':    'slideLayout11.xml',
    'left_arrow':         'slideLayout12.xml',
    'green_left_arrow':   'slideLayout13.xml',
    'arrow_one_third':    'slideLayout14.xml',
    'green_arrow_one_third': 'slideLayout15.xml',
    'arrow_half':         'slideLayout16.xml',
    'green_arrow_half':   'slideLayout17.xml',
    'arrow_two_third':    'slideLayout18.xml',
    'green_arrow_two_third': 'slideLayout19.xml',
    'big_statement_green':'slideLayout20.xml',
    'big_statement_icon': 'slideLayout21.xml',
    'quote':              'slideLayout22.xml',
    'special_gray':       'slideLayout23.xml',
    'blank':              'slideLayout24.xml',
    'blank_green':        'slideLayout25.xml',
    'disclaimer':         'slideLayout26.xml',
    'end':                'slideLayout27.xml',
    # Detail variants (D. prefix) -- every standard layout has a detail counterpart
    # Detail layouts have white/light backgrounds and smaller titles for data-rich slides
    'd_title':              'slideLayout29.xml',
    'd_title_only':         'slideLayout30.xml',
    'd_title_and_text':     'slideLayout31.xml',
    'd_gray_slice':         'slideLayout32.xml',
    'd_section_header_box': 'slideLayout33.xml',
    'd_section_header_line':'slideLayout34.xml',
    'd_white_one_third':    'slideLayout35.xml',
    'd_green_highlight':    'slideLayout36.xml',
    'd_four_column_green':  'slideLayout37.xml',
    'd_green_one_third':    'slideLayout38.xml',
    'd_green_half':         'slideLayout39.xml',
    'd_green_two_third':    'slideLayout40.xml',
    'd_left_arrow':         'slideLayout41.xml',
    'd_green_left_arrow':   'slideLayout42.xml',
    'd_arrow_one_third':    'slideLayout43.xml',
    'd_green_arrow_one_third': 'slideLayout44.xml',
    'd_arrow_half':         'slideLayout45.xml',
    'd_green_arrow_half':   'slideLayout46.xml',
    'd_arrow_two_third':    'slideLayout47.xml',
    'd_green_arrow_two_third': 'slideLayout48.xml',
    'd_big_statement_green':'slideLayout49.xml',
    'd_big_statement_icon': 'slideLayout50.xml',
    'd_quote':              'slideLayout51.xml',
    'd_special_gray':       'slideLayout52.xml',
    'd_blank':              'slideLayout55.xml',
    'd_blank_green':        'slideLayout54.xml',
    'd_disclaimer':         'slideLayout56.xml',
    'd_end':                'slideLayout57.xml',
}


# ============================================================
# ICON SEARCH (v6.5 -- ported from deckster-digital token search)
# ============================================================

_ICON_STOP_WORDS = frozenset({
    'a','an','the','and','or','but','in','on','at','to','for','of','with',
    'by','from','is','it','as','be','was','are','been','being','has','have',
    'had','do','does','did','will','would','could','should','may','might',
    'can','shall','not','no','nor','so','if','then','than','too','very',
    'just','about','above','after','again','all','also','am','any','because',
    'before','below','between','both','during','each','few','further','get',
    'got','he','her','here','him','his','how','i','into','its','let','me',
    'more','most','my','new','now','off','old','only','other','our','out',
    'over','own','re','same','she','some','still','such','take','their',
    'them','there','these','they','this','those','through','under','up',
    'us','we','what','when','where','which','while','who','whom','why',
    'you','your',
})

_ICON_QUERY_SYNONYMS = {
    'logistics': ('supply', 'chain', 'distribution', 'warehouse', 'container', 'truck', 'shipping', 'delivery'),
    'supply': ('logistics', 'chain', 'distribution', 'warehouse', 'container'),
    'distribution': ('logistics', 'supply', 'network', 'warehouse', 'delivery'),
    'finance': ('bank', 'money', 'dollar', 'cash', 'growth', 'chart'),
    'money': ('finance', 'bank', 'dollar', 'cash'),
    'people': ('team', 'person', 'group', 'meeting', 'community'),
    'team': ('people', 'group', 'collaboration', 'community'),
    'person': ('people', 'team', 'individual'),
    'partnership': ('partner', 'collaboration', 'handshake', 'team'),
    'partner': ('partnership', 'collaboration', 'handshake', 'team'),
    'global': ('globe', 'world', 'earth', 'map', 'network'),
    'globe': ('global', 'world', 'earth', 'map'),
    'world': ('global', 'globe', 'earth', 'map'),
    'growth': ('chart', 'arrow', 'target', 'digital'),
    'strategy': ('target', 'roadmap', 'direction', 'growth'),
    'target': ('strategy', 'goal', 'aim', 'bullseye'),
    'medical': ('health', 'hospital', 'heart', 'pharma', 'pill'),
    'health': ('medical', 'hospital', 'heart', 'pharma', 'pill'),
    'pharma': ('medical', 'health', 'pill', 'hospital'),
    'network': ('connection', 'map', 'global', 'organization'),
    'org': ('organization', 'hierarchy', 'people', 'team', 'building'),
    'organization': ('org', 'hierarchy', 'team', 'people'),
    'building': ('office', 'bank', 'factory', 'organization'),
    'location': ('map', 'pin', 'global', 'place'),
    'risk': ('shield', 'alert', 'warning', 'decision', 'check'),
    'decision': ('check', 'target', 'crossroads', 'direction'),
    'award': ('trophy', 'star', 'ribbon', 'best'),
    'trophy': ('award', 'star', 'ribbon', 'best'),
    'idea': ('light', 'bulb', 'innovation', 'spark'),
    'innovation': ('idea', 'light', 'bulb', 'digital'),
}


def _split_identifier_tokens(text):
    """Split human text and CamelCase identifiers into searchable tokens."""
    parts = []
    for chunk in re.split(r'[\s/_-]+', text):
        if not chunk:
            continue
        camel_parts = re.findall(r'[A-Z]+(?=[A-Z][a-z]|[0-9]|$)|[A-Z]?[a-z]+|[0-9]+', chunk)
        if camel_parts:
            parts.extend(camel_parts)
        else:
            parts.append(chunk)
    return parts


def _tokenize(text):
    """Normalize and tokenize text for icon search. Filters stop words."""
    normalized = (
        text.replace('bcgIcons_', '')
        .replace('svg_', '')
        .replace('_bug', '')
    )
    tokens = [part.lower() for part in _split_identifier_tokens(normalized)]
    return [t for t in tokens if t not in _ICON_STOP_WORDS and len(t) > 1]


def _expand_icon_query_tokens(tokens):
    """Expand common business-domain tokens into likely icon-name synonyms."""
    expanded = []
    seen = set()
    for token in tokens:
        if token not in seen:
            expanded.append(token)
            seen.add(token)
        for synonym in _ICON_QUERY_SYNONYMS.get(token, ()):
            if synonym not in seen:
                expanded.append(synonym)
                seen.add(synonym)
    return expanded


def _display_icon_name(icon_key, icon_type='icon'):
    """Convert bundled icon keys into display names used in docs and search output."""
    name = icon_key.replace('bcgIcons_', '').replace('svg_', '')
    if icon_type == 'bug' and name.endswith('_bug'):
        name = name[:-4]
    return name


def _normalize_icon_lookup(value):
    """Normalize icon names/queries for exact lookup."""
    return re.sub(r'[^a-z0-9]+', '', str(value or '').lower())


# ============================================================
# THEME CONFIGURATION (v6.5 -- ee4p client template support)
# ============================================================

def apply_theme(config):
    """Override module-level constants from a template.json config.

    Called by BCGDeck.__init__ when theme_config is provided.
    Modifies global FONT, COLORS, LAYOUTS, and positioning constants.

    Args:
        config: dict from template.json (output of ingest_ee4p.py)
    """
    global FONT, COLORS, CONTENT_START_Y, DETAIL_CONTENT_START_Y
    global LIGHT_PALETTE, DARK_PALETTE
    global CONTENT_END_Y, SOURCE_Y, CONTENT_LEFT, CONTENT_WIDTH, LAYOUTS
    global _LAYOUT_BACKGROUNDS, _ACTIVE_THEME_CONFIG

    config = normalize_theme_config(config)

    # Store config for content_bounds() lookups
    _ACTIVE_THEME_CONFIG = config

    # Font
    if 'font' in config and config['font']:
        FONT = config['font']

    # Color aliases -> COLORS dict
    aliases = config.get('color_aliases', {})
    alias_to_color = {
        'GREEN': 'BCG_GREEN', 'DARK_GREEN': 'DARK_GREEN', 'DEEP_GREEN': 'DEEP_GREEN',
        'DARK_TEXT': 'DARK_TEXT', 'LIGHT_BG': 'LIGHT_BG', 'WHITE': 'WHITE',
        'MED_GRAY': 'MED_GRAY',
    }
    for alias_key, color_key in alias_to_color.items():
        if alias_key in aliases:
            COLORS[color_key] = aliases[alias_key]

    # Map scheme colors to semantic COLORS entries
    scheme = config.get('scheme_colors', {})
    if scheme:
        for sc_key, colors_key in {
            'accent1': 'ACCENT1', 'accent2': 'ACCENT2', 'accent3': 'ACCENT3',
            'accent4': 'ACCENT4', 'accent5': 'ACCENT5', 'accent6': 'ACCENT6',
            'lt1': 'SURFACE', 'lt2': 'SURFACE_ALT',
            'dk1': 'ON_SURFACE', 'dk2': 'CONTRAST_BG',
        }.items():
            if sc_key in scheme:
                COLORS[colors_key] = scheme[sc_key].lstrip('#')
        if 'lt2' in scheme:
            COLORS['OUTLINE'] = scheme['lt2'].lstrip('#')
        if 'lt1' in scheme:
            COLORS['CONTRAST_FG'] = scheme['lt1'].lstrip('#')

    # Wire title/body styles from extraction (overrides color aliases)
    title_style = config.get('title_style', {})
    if title_style.get('fillColor'):
        COLORS['TITLE_COLOR'] = title_style['fillColor'].lstrip('#')
    if title_style.get('fontName'):
        FONT = title_style['fontName']

    body_style = config.get('body_style', {})
    if body_style.get('fillColor'):
        COLORS['DARK_TEXT'] = body_style['fillColor'].lstrip('#')
    if body_style.get('fontName') and not config.get('font'):
        FONT = body_style['fontName']

    # Remap BCG-specific color keys to client accents so nothing leaks
    # These keys are used by methods like add_chevron_flow, add_label defaults, etc.
    # Skip for BCG default — the hardcoded COLORS already match the scheme_colors,
    # and remapping would scramble them (dk2=29BA74 would overwrite DEEP_GREEN, etc.)
    is_bcg_default = config.get('name', '') == 'bcg_default' or config.get('source', '') == 'built-in'
    if scheme and not is_bcg_default:
        # DARK_GREEN/DEEP_GREEN -> use darkest accent or dk2
        dk2 = scheme.get('dk2', '').lstrip('#')
        a1 = scheme.get('accent1', '').lstrip('#')
        a2 = scheme.get('accent2', '').lstrip('#')
        a3 = scheme.get('accent3', '').lstrip('#')
        a5 = scheme.get('accent5', '').lstrip('#')
        a6 = scheme.get('accent6', '').lstrip('#')

        if dk2:
            COLORS['DEEP_GREEN'] = dk2
        if a1:
            COLORS['DARK_GREEN'] = a1
            COLORS['BCG_GREEN'] = a1  # Remap primary green to client accent1

        # TEAL, LIME, NAVY -> map to remaining accents
        if a5:
            COLORS['TEAL'] = a5
        if a6:
            COLORS['LIME'] = a6
        if dk2:
            COLORS['NAVY'] = dk2

        # NEGATIVE/STAMP_RED -> use accent2 (typically warm/red)
        if a2:
            COLORS['NEGATIVE'] = a2
            COLORS['STAMP_RED'] = a2

        # PAGE_NUM -> muted gray from theme
        lt2 = scheme.get('lt2', '').lstrip('#')
        if lt2:
            COLORS['PAGE_NUM'] = COLORS.get('MED_GRAY', 'B0B0B0')

        # LIGHT_TEXT -> match surface alt
        if lt2:
            COLORS['LIGHT_TEXT'] = lt2

    # Store slide background color from scheme (bg1 = primary slide bg)
    if scheme:
        bg1 = scheme.get('bg1', scheme.get('lt1', '')).lstrip('#')
        if bg1:
            COLORS['SLIDE_BG'] = bg1

    # Store layout background data for layout-aware text color selection
    layout_geom = config.get('layout_geometry', {})
    if layout_geom:
        _LAYOUT_BACKGROUNDS = {}
        for slug, info in layout_geom.items():
            # Priority: accent region > content_zone.content_bg > layout background > slot bg
            accents = info.get('accent_regions', [])
            cz = info.get('content_zone', {})
            layout_bg = (info.get('background') or '').lstrip('#')
            slots = info.get('slots', [])

            if accents:
                _LAYOUT_BACKGROUNDS[slug] = {
                    'accent_fill': (accents[0].get('fill') or '').lstrip('#'),
                    'palette_role': accents[0].get('paletteRole', 'light'),
                }
            elif cz and cz.get('content_bg'):
                _LAYOUT_BACKGROUNDS[slug] = {
                    'content_bg': cz['content_bg'].lstrip('#'),
                    'palette_role': cz.get('content_palette_role', 'light'),
                }
            elif layout_bg:
                # Determine palette role from luminance
                _bg_rgb = _hex_to_rgb(layout_bg)
                bg_role = 'dark' if (_bg_rgb and _relative_luminance(_bg_rgb) < 0.4) else 'light'
                _LAYOUT_BACKGROUNDS[slug] = {
                    'layout_bg': layout_bg,
                    'palette_role': bg_role,
                }
            elif slots:
                _LAYOUT_BACKGROUNDS[slug] = {
                    'slot_bg': slots[0].get('background', '').lstrip('#'),
                    'palette_role': slots[0].get('paletteRole', 'light'),
                }

    # Build light/dark palettes from resolved COLORS
    LIGHT_PALETTE['header'] = COLORS.get('TITLE_COLOR', COLORS.get('ACCENT1', COLORS['BCG_GREEN']))
    LIGHT_PALETTE['text'] = COLORS['DARK_TEXT']
    LIGHT_PALETTE['secondary'] = COLORS.get('MED_GRAY', '6E6F73')

    DARK_PALETTE['header'] = COLORS.get('CONTRAST_FG', 'FFFFFF')
    DARK_PALETTE['text'] = COLORS.get('CONTRAST_FG', 'FFFFFF')
    DARK_PALETTE['secondary'] = COLORS.get('SURFACE_ALT', 'F2F2F2')

    # Positioning from fly_zone
    fz = config.get('fly_zone', {})
    if fz:
        if 'content_left' in fz:
            CONTENT_LEFT = fz['content_left']
        if 'content_width' in fz:
            CONTENT_WIDTH = fz['content_width']
        if 'content_top' in fz:
            CONTENT_START_Y = fz['content_top']
        if 'content_bottom' in fz:
            CONTENT_END_Y = fz['content_bottom']
            SOURCE_Y = round(CONTENT_END_Y + 0.24, 2)

    # Layout XML filenames from xml_layout_map + semantic_layout_map
    xml_map = config.get('xml_layout_map', {})
    semantic = config.get('semantic_layout_map', {})
    if xml_map and semantic:
        for sem_key, layout_key in SEMANTIC_TO_RUNTIME_LAYOUT.items():
            layout_name = semantic.get(sem_key)
            if layout_name and layout_name in xml_map:
                LAYOUTS[layout_key] = xml_map[layout_name]

        # Detail layout map
        detail_map = config.get('detail_layout_map', {})
        if detail_map:
            for base_name, full_name in detail_map.items():
                if full_name in xml_map:
                    normalized = re.sub(r'[^a-z0-9]+', '_', base_name.lower()).strip('_')
                    d_key = 'd_' + normalized
                    LAYOUTS[d_key] = xml_map[full_name]

        # Fallback: for d_* layouts not explicitly set by the client theme,
        # point them to the non-detail equivalent. This ensures `detail=True`
        # uses a valid client layout instead of a stale BCG slideLayout.
        theme_set_layouts = set()  # track what this theme explicitly mapped
        for sem_key, layout_name in semantic.items():
            if layout_name in xml_map:
                theme_set_layouts.add(xml_map[layout_name])
        if detail_map:
            for bn, fn in detail_map.items():
                if fn in xml_map:
                    theme_set_layouts.add(xml_map[fn])

        for key in list(LAYOUTS.keys()):
            if key.startswith('d_') and LAYOUTS[key] not in theme_set_layouts:
                base_key = key[2:]
                if base_key in LAYOUTS and LAYOUTS[base_key] in theme_set_layouts:
                    LAYOUTS[key] = LAYOUTS[base_key]


# Auto-generated standard->detail layout lookup
_DETAIL_MAP = {k: f'd_{k}' for k in LAYOUTS if not k.startswith('d_') and f'd_{k}' in LAYOUTS}

# Picture placeholder info for layouts with dedicated image zones.
# idx = placeholder index in the layout; x/y/w/h = position in inches.
_PICTURE_PLACEHOLDERS = {
    'green_half':        {'idx': 14, 'x': 6.66, 'y': 0.00, 'w': 6.67, 'h': 7.50},
    'green_two_third':   {'idx': 11, 'x': 8.55, 'y': 0.00, 'w': 4.78, 'h': 7.50},
    'd_green_half':      {'idx': 14, 'x': 6.66, 'y': 0.00, 'w': 6.67, 'h': 7.50},
    'd_green_two_third': {'idx': 11, 'x': 8.55, 'y': 0.00, 'w': 4.78, 'h': 7.50},
}

# ============================================================
# IMAGE CATALOG (tag-based stock image selection)
# ============================================================

_IMAGE_DIR = find_skill_root() / 'assets' / 'images'
_IMAGE_BUNDLE_PATH = _IMAGE_DIR / 'images_bundle.json'
_IMAGE_BUNDLE = None  # lazy-loaded raw bundle
_IMAGE_LOOKUP = None  # filename -> bundled image metadata
_IMAGE_CATALOG = None  # lazy-loaded selection catalog
_IMAGE_CONTENT_TYPES = {
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.emf': 'image/x-emf',
    '.gif': 'image/gif',
    '.svg': 'image/svg+xml',
    '.tiff': 'image/tiff',
    '.tif': 'image/tiff',
    '.bmp': 'image/bmp',
}


def _normalize_token(value):
    """Lowercase alphanumeric only -- for fuzzy tag matching."""
    return re.sub(r'[^a-z0-9]+', '', value.strip().lower())


def _image_content_type(ext):
    """Return the OOXML content type for a file extension."""
    return _IMAGE_CONTENT_TYPES.get(ext.lower(), 'image/png')


def _load_image_catalog():
    """Load and cache the bundled stock-image catalog."""
    global _IMAGE_BUNDLE, _IMAGE_LOOKUP, _IMAGE_CATALOG
    if _IMAGE_CATALOG is not None:
        return _IMAGE_CATALOG

    if not _IMAGE_BUNDLE_PATH.exists():
        _IMAGE_BUNDLE = {}
        _IMAGE_LOOKUP = {}
        _IMAGE_CATALOG = []
        return _IMAGE_CATALOG

    try:
        raw = json.loads(_IMAGE_BUNDLE_PATH.read_text())
        images = raw.get('images', {})
        if not isinstance(images, dict):
            _IMAGE_BUNDLE = {}
            _IMAGE_LOOKUP = {}
            _IMAGE_CATALOG = []
            return _IMAGE_CATALOG

        catalog = []
        lookup = {}
        for key, item in images.items():
            if not isinstance(item, dict):
                continue
            filename = str(item.get('filename', key)).strip()
            if not filename:
                continue
            tags = item.get('tags', [])
            if not isinstance(tags, list):
                tags = []
            desc = str(item.get('description', ''))
            size_kb = float(item.get('size_kb', 0.0))
            bundle_item = dict(item)
            bundle_item['filename'] = filename
            ext = Path(filename).suffix.lower()
            bundle_item.setdefault('mime_type', _image_content_type(ext))
            lookup[filename] = bundle_item
            catalog.append({
                'filename': filename,
                'path': _IMAGE_DIR / filename,
                'tags': [str(t) for t in tags],
                'description': desc,
                'size_kb': size_kb,
                'width': item.get('width'),
                'height': item.get('height'),
                'mime_type': str(bundle_item.get('mime_type')),
                'bundle_key': filename,
            })
        _IMAGE_BUNDLE = raw
        _IMAGE_LOOKUP = lookup
        _IMAGE_CATALOG = catalog
        return _IMAGE_CATALOG
    except Exception:
        _IMAGE_BUNDLE = {}
        _IMAGE_LOOKUP = {}
        _IMAGE_CATALOG = []
        return _IMAGE_CATALOG


def _lookup_bundled_image(image_ref):
    """Resolve an image reference to a bundled image entry, if present."""
    _load_image_catalog()
    if not _IMAGE_LOOKUP:
        return None

    ref = str(image_ref).strip()
    candidates = []
    if ref:
        candidates.append(ref)
        name = Path(ref).name
        if name and name != ref:
            candidates.append(name)

    for candidate in candidates:
        item = _IMAGE_LOOKUP.get(candidate)
        if item:
            return item
    return None


def _resolve_image_source(image_ref):
    """Resolve an image path or bundled stock-image ref to bytes + metadata."""
    path = Path(image_ref)
    if path.exists():
        image_bytes = path.read_bytes()
        ext = path.suffix.lower()
        return {
            'bytes': image_bytes,
            'ext': ext,
            'filename': path.name,
            'content_type': _image_content_type(ext),
        }

    bundled = _lookup_bundled_image(image_ref)
    if bundled:
        ext = Path(bundled['filename']).suffix.lower()
        return {
            'bytes': zlib.decompress(base64.b64decode(bundled['z'])),
            'ext': ext,
            'filename': bundled['filename'],
            'content_type': bundled.get('mime_type') or _image_content_type(ext),
        }

    raise FileNotFoundError(f"Image not found: {image_ref}")


def _select_image_from_tags(tags, match_mode='any'):
    """Find the best matching stock image by tags.

    Args:
        tags: List of tag strings (e.g. ['logistics', 'warehouse'])
        match_mode: 'any' (at least one tag) or 'all' (every tag must match)

    Returns:
        Asset ref path for the best image, or None if no match.
    """
    query = [_normalize_token(t) for t in tags if _normalize_token(t)]
    if not query:
        return None

    mode = match_mode.strip().lower() if isinstance(match_mode, str) else 'any'
    if mode not in ('any', 'all'):
        mode = 'any'

    best = None
    best_path = None

    for image in _load_image_catalog():
        # Build vocabulary from tags + description words
        vocab = set()
        for tag in image['tags']:
            token = _normalize_token(str(tag))
            if token:
                vocab.add(token)
        for word in re.findall(r'[a-z0-9]+', image['description'].lower()):
            token = _normalize_token(word)
            if token:
                vocab.add(token)

        matches = sum(1 for term in query if term in vocab)
        if mode == 'all' and matches < len(query):
            continue
        if matches == 0:
            continue

        score = (matches, image['size_kb'], image['filename'])
        if best is None or score > best:
            best = score
            best_path = image['path']

    return best_path


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def emu(inches):
    """Convert inches to EMU."""
    return str(int(inches * 914400))

def pt_emu(pt):
    """Convert points to EMU (for font sizes in hundredths of a point)."""
    return str(int(pt * 100))

def _escape_xml(text):
    """Escape text for XML inclusion."""
    return (text
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
        .replace('"', '&quot;')
        .replace("'", '&apos;'))

def _valign_map(v):
    """Map valign string to OOXML anchor value."""
    return {'top': 't', 'middle': 'ctr', 'bottom': 'b'}.get(v, 't')

def _shadow_xml():
    """Subtle outer shadow matching reference deck: blur=4pt, dist=2pt, dir=135deg, alpha=8%."""
    return (
        '<a:effectLst>'
        '<a:outerShdw blurRad="50800" dist="25400" dir="8100000" algn="tl" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="8000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )


_VALID_HEX = re.compile(r'^[0-9A-Fa-f]{6}$')

def _sanitize_color(val):
    """Normalize a color value to exactly 6 hex chars for OOXML srgbClr val.

    Handles common mistakes:
      - '#RRGGBB' -> strip leading '#'
      - 'RRGGBBAA' (8-char with alpha) -> truncate to 'RRGGBB'
      - 'RGB' (3-char CSS shorthand) -> expand to 'RRGGBB'
    Raises ValueError if the result is not valid hex.
    """
    if val is None:
        return None
    s = str(val).strip().lstrip('#')
    # 8-char RGBA -> drop alpha
    if len(s) == 8 and all(c in '0123456789ABCDEFabcdef' for c in s):
        s = s[:6]
    # 3-char CSS shorthand -> expand
    if len(s) == 3 and all(c in '0123456789ABCDEFabcdef' for c in s):
        s = ''.join(c * 2 for c in s)
    if not _VALID_HEX.match(s):
        raise ValueError(
            f'Invalid color "{val}": must be 6-character hex (RRGGBB), '
            f'e.g. "29BA74". Got {len(s)} chars after cleanup.'
        )
    return s.upper()


def _run_xml(text, sz=12, color=None, bold=False, italic=False, underline=False,
             strikethrough=False, font=None, highlight=None, superscript=False, subscript=False):
    color = color or COLORS['DARK_TEXT']
    """Generate a single <a:r> XML element with formatting.

    Args:
        text: Plain text content for this run.
        sz: Font size in points.
        color: Hex color (6 chars, no #).
        bold: Bold text.
        italic: Italic text.
        underline: Underline text.
        strikethrough: Strikethrough text.
        font: Font override (default: FONT constant = Trebuchet MS).
        highlight: Background highlight hex color (uses a:highlight).
        superscript: Superscript (baseline shift +30000).
        subscript: Subscript (baseline shift -25000).
    """
    color = _sanitize_color(color)
    highlight = _sanitize_color(highlight)
    escaped = _escape_xml(text)
    attrs = f'lang="en-US" sz="{pt_emu(sz)}" dirty="0"'
    if bold:
        attrs += ' b="1"'
    if italic:
        attrs += ' i="1"'
    if underline:
        attrs += ' u="sng"'
    if strikethrough:
        attrs += ' strike="sngStrike"'
    if superscript:
        attrs += ' baseline="30000"'
    elif subscript:
        attrs += ' baseline="-25000"'

    fill_xml = f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    hl_xml = f'<a:highlight><a:srgbClr val="{highlight}"/></a:highlight>' if highlight else ''
    font_xml = f'<a:latin typeface="{font or FONT}"/>'

    return (f'<a:r><a:rPr {attrs}>{fill_xml}{hl_xml}{font_xml}</a:rPr>'
            f'<a:t>{escaped}</a:t></a:r>')


def _para_xml(runs, align='left', bullet=False, level=0, spc_after=0):
    """Generate a single <a:p> XML element from a list of runs.

    Args:
        runs: List of run XML strings (from _run_xml).
        align: Paragraph alignment: 'left', 'center', 'right'.
        bullet: If True, add bullet character.
        level: Indentation level for bullets (0 = top level).
        spc_after: Space after paragraph in points (0=default, 6=comfortable, 10=spacious).
    """
    algn_map = {'left': 'l', 'center': 'ctr', 'right': 'r'}
    algn_val = algn_map.get(align, 'l')
    spc_xml = f'<a:spcAft><a:spcPts val="{int(spc_after * 100)}"/></a:spcAft>' if spc_after else ''

    if bullet:
        mar = 228600 * (level + 1)
        indent = -228600
        # OOXML schema requires spcAft BEFORE buChar inside <a:pPr>
        return (f'<a:p><a:pPr marL="{mar}" indent="{indent}" algn="{algn_val}">'
                f'{spc_xml}<a:buChar char="-"/></a:pPr>{"".join(runs)}</a:p>')
    elif align != 'left' or spc_xml:
        return f'<a:p><a:pPr algn="{algn_val}">{spc_xml}</a:pPr>' + ''.join(runs) + '</a:p>'
    else:
        return '<a:p>' + ''.join(runs) + '</a:p>'


def _resolve_runs(text_or_runs, sz=12, color=None, bold=False):
    color = color or COLORS['DARK_TEXT']
    """Convert text input to a list of run XML strings.

    Accepts:
        str: Plain text -- creates a single run with the default formatting.
        list: Rich text -- list of tuples/dicts defining multiple runs.
              Each element can be:
                - (text, {props})  -- text with property overrides
                - {'t': text, ...props}  -- dict with 't' key for text
              Props override the defaults: sz, color, bold, italic,
              underline, strikethrough, font, highlight, superscript, subscript.
    """
    if isinstance(text_or_runs, str):
        return [_run_xml(text_or_runs, sz=sz, color=color, bold=bold)]

    runs = []
    for item in text_or_runs:
        if isinstance(item, dict):
            t = item.get('t', '')
            props = {k: v for k, v in item.items() if k != 't'}
        elif isinstance(item, (tuple, list)) and len(item) == 2:
            t, props = item[0], item[1] if isinstance(item[1], dict) else {}
        else:
            t, props = str(item), {}

        runs.append(_run_xml(
            t,
            sz=props.get('sz', sz),
            color=props.get('color', color),
            bold=props.get('bold', bold),
            italic=props.get('italic', False),
            underline=props.get('underline', False),
            strikethrough=props.get('strikethrough', False),
            font=props.get('font', None),
            highlight=props.get('highlight', None),
            superscript=props.get('superscript', False),
            subscript=props.get('subscript', False),
        ))
    return runs


def _make_text_runs(text, sz=12, color=None, bold=False, align='left'):
    color = color or COLORS['DARK_TEXT']
    """Create XML paragraphs from text input.

    text can be:
        str: Plain text. Newlines become separate paragraphs.
        list of runs: Rich text -- all runs in one paragraph.
        list of lists: Multiple paragraphs, each with its own runs.
            e.g. [  [('Revenue: ', {}), ('$2.8B', {'bold': True, 'color': '29BA74'})],
                    [('Growth: ', {}), ('+23%', {'bold': True, 'color': '29BA74'})]  ]
    """
    # Plain string -- split on newlines
    if isinstance(text, str):
        paragraphs = text.split('\n')
        parts = []
        for para in paragraphs:
            runs = _resolve_runs(para, sz=sz, color=color, bold=bold)
            parts.append(_para_xml(runs, align=align))
        return '\n'.join(parts)

    # List input -- check if it's a list of paragraphs or a single paragraph of runs
    if text and isinstance(text[0], list):
        # List of paragraphs, each paragraph is a list of runs
        parts = []
        for para_runs in text:
            runs = _resolve_runs(para_runs, sz=sz, color=color, bold=bold)
            parts.append(_para_xml(runs, align=align))
        return '\n'.join(parts)
    else:
        # Single paragraph of runs
        runs = _resolve_runs(text, sz=sz, color=color, bold=bold)
        return _para_xml(runs, align=align)


def _make_bullet_runs(items, sz=12, color=None, align='left', spc_after=0):
    color = color or COLORS['DARK_TEXT']
    """Create XML bullet list paragraphs.

    Each item can be:
        str: Plain bullet text.
        list of runs: Rich text bullet -- mixed formatting within one bullet.
            e.g. [('Revenue grew ', {}), ('+23%', {'bold': True, 'color': '29BA74'})]
        dict with 'runs' key: Rich text with indent level.
            e.g. {'runs': [('Key point', {'bold': True})], 'level': 1}
    spc_after: Space after each bullet in points (0=default, 6=comfortable, 10=spacious).
    """
    parts = []
    for item in items:
        if isinstance(item, dict) and 'runs' in item:
            level = item.get('level', 0)
            runs = _resolve_runs(item['runs'], sz=sz, color=color)
            parts.append(_para_xml(runs, align=align, bullet=True, level=level, spc_after=spc_after))
        elif isinstance(item, list):
            runs = _resolve_runs(item, sz=sz, color=color)
            parts.append(_para_xml(runs, align=align, bullet=True, spc_after=spc_after))
        else:
            runs = _resolve_runs(str(item), sz=sz, color=color)
            parts.append(_para_xml(runs, align=align, bullet=True, spc_after=spc_after))
    return '\n'.join(parts)


# ============================================================
# BCGDECK CLASS
# ============================================================

def check_setup(script_path=None):
    """Validate that the BCG slide generator environment is correctly set up.

    Checks that all required assets (template, icons, images, dependencies)
    are accessible from the given script location. Call this before building
    any deck to catch path/setup issues early.

    Args:
        script_path: Path to bcg_template.py. Defaults to this file's location.

    Returns:
        True if all checks pass, False otherwise. Prints diagnostic output.
    """
    base = find_skill_root(script_path)
    ok = True
    checks = []

    # 1. Template .pptx (lives in styles/templates/bcg_default/)
    tpl = base / 'styles' / 'templates' / 'bcg_default' / 'master.pptx'
    if not tpl.exists():
        # Fallback: legacy location in assets/
        tpl = base / 'assets' / 'BCG_Master_16-9_Default.pptx'
    if tpl.exists():
        checks.append(('Template .pptx', 'OK', str(tpl)))
    else:
        checks.append(('Template .pptx', 'MISSING', 'styles/templates/bcg_default/master.pptx'))
        ok = False

    # 2. Icons bundle
    icons_bundle = base / 'assets' / 'icons' / 'icons_bundle.json'
    if icons_bundle.exists():
        try:
            with open(icons_bundle) as f:
                n_icons = len(json.load(f).get('icons', {}))
            checks.append(('Icons library', 'OK', f'{n_icons} icons'))
        except (json.JSONDecodeError, KeyError, TypeError) as e:
            checks.append(('Icons library', 'FAIL', f'corrupted icons_bundle.json: {e}'))
            ok = False
    else:
        checks.append(('Icons library', 'MISSING', str(icons_bundle)))
        ok = False

    # 3. Images bundle
    imgs_bundle = base / 'assets' / 'images' / 'images_bundle.json'
    if imgs_bundle.exists():
        try:
            with open(imgs_bundle) as f:
                n_imgs = len(json.load(f).get('images', {}))
            checks.append(('Stock images', 'OK', f'{n_imgs} images'))
        except (json.JSONDecodeError, KeyError, TypeError) as e:
            checks.append(('Stock images', 'FAIL', f'corrupted images_bundle.json: {e}'))
            ok = False
    else:
        checks.append(('Stock images', 'MISSING', str(imgs_bundle)))
        ok = False

    # 4. pptx_utils.py co-located
    utils = Path(script_path or __file__).resolve().parent / 'pptx_utils.py'
    if utils.exists():
        checks.append(('pptx_utils.py', 'OK', str(utils)))
    else:
        checks.append(('pptx_utils.py', 'MISSING', str(utils)))
        ok = False

    # 5. python-pptx (optional, for charts)
    try:
        import pptx  # noqa: F401
        checks.append(('python-pptx', 'OK', 'installed (charts available)'))
    except ImportError:
        checks.append(('python-pptx', 'WARN', 'not installed (charts will fail -- pip install python-pptx)'))

    # Print report
    print(f"\nBCG Slide Generator -- Setup Check")
    print(f"Skill directory: {base}\n")
    for name, status, detail in checks:
        icon = '  OK' if status == 'OK' else 'WARN' if status == 'WARN' else 'FAIL'
        print(f"  [{icon}] {name}: {detail}")

    if ok:
        print(f"\nAll checks passed. Ready to build decks.")
    else:
        print(f"\nSetup issues found. The skill directory should contain:")
        print(f"  {base}/")
        print(f"  |---- scripts/bcg_template.py")
        print(f"  |---- scripts/pptx_utils.py")
        print(f"  |---- styles/templates/bcg_default/master.pptx")
        print(f"  \---- assets/")
        print(f"      |---- icons/icons_bundle.json")
        print(f"      \---- images/images_bundle.json")
        print(f"\n  If you have multiple skills installed, ensure the import path")
        print(f"  points to the deckster-slide-generator skill specifically.")

    return ok


class BCGDeck:
    """Build a BCG presentation from the official master template."""

    def __init__(self, template_path=None, theme_config=None, unpack_template=True):
        """Initialize by unpacking the BCG master template.

        Args:
            template_path: Path to master .pptx. If None, uses BCG default
                          or the one from theme_config.
            theme_config: Path to template.json or dict from ingest_ee4p.
                         When provided, overrides colors, fonts, layouts,
                         and positioning to match the client template.
            unpack_template: When False, skips unpacking the master .pptx and
                            only initializes theme/search state. Use this for
                            fast preflight, title validation, and icon resolution.
        """
        # Load template config — either explicit client theme or bundled bcg_default
        if theme_config is not None:
            if isinstance(theme_config, (str, Path)):
                with open(theme_config) as _f:
                    theme_config = json.load(_f)
            theme_config = normalize_theme_config(theme_config)
        else:
            # Auto-load bcg_default template.json for unified code path
            try:
                from template_registry import get_template
                theme_config = get_template(name='bcg_default')
            except (ImportError, Exception):
                # Fallback: try to find template.json directly
                skill_dir = find_skill_root()
                bcg_config = skill_dir / 'styles' / 'templates' / 'bcg_default' / 'template.json'
                if not bcg_config.exists():
                    bcg_config = skill_dir / 'templates' / 'bcg_default' / 'template.json'
                if bcg_config.exists():
                    with open(bcg_config) as _f:
                        theme_config = json.load(_f)
                    theme_config['_template_dir'] = str(bcg_config.parent)
                    theme_config['_master_pptx'] = str(bcg_config.parent / 'master.pptx')
                    theme_config = normalize_theme_config(theme_config)

        # Apply theme (runs for both BCG default and client templates)
        if theme_config is not None:
            apply_theme(theme_config)
            # Use template's master pptx if no explicit template_path
            if template_path is None and '_master_pptx' in theme_config:
                template_path = theme_config['_master_pptx']
            elif template_path is None and 'template_file' in theme_config:
                config_dir = Path(theme_config.get('_template_dir', '.'))
                template_path = config_dir / theme_config['template_file']

        if template_path is None:
            # Last resort fallback — try bundled template, then legacy assets
            skill_dir = find_skill_root()
            template_path = skill_dir / 'styles' / 'templates' / 'bcg_default' / 'master.pptx'
            if not template_path.exists():
                template_path = skill_dir / 'assets' / 'BCG_Master_16-9_Default.pptx'
            if not template_path.exists():
                template_path = Path('BCG_Master_16-9_Default.pptx')

        if not Path(template_path).exists():
            raise FileNotFoundError(
                f"BCG master template not found at {template_path}. "
                f"Provide the path to BCG_Master_16-9_Default.pptx."
            )

        self.template_path = str(template_path)
        self._shape_id_counter = 100
        self._slide_count = 0
        self._chart_specs = []  # Deferred chart injection via python-pptx
        self._slide_layouts = {}  # slide_path -> layout_key mapping
        self._slide_contracts = {}  # slide_path -> runtime slide contract
        self._variant_preferences = {}  # pattern -> variant name overrides for render_pattern()
        self._theme_config = theme_config  # Stored for pattern_variants.py fallback resolution

        self.work_dir = make_work_dir('bcg_deck_') if unpack_template else None

        if unpack_template:
            # Unpack template
            unpack(self.template_path, str(self.work_dir))

            # Clear existing placeholder slides
            self._clear_slides()

    def set_pattern_variant(self, pattern, variant):
        """Override the default render_pattern() variant for this deck."""
        if not pattern or not variant:
            raise ValueError("Both pattern and variant are required")
        self._variant_preferences[str(pattern)] = str(variant)

    def set_pattern_variants(self, mapping):
        """Override multiple render_pattern() defaults for this deck."""
        if mapping is None:
            return
        for pattern, variant in dict(mapping).items():
            self.set_pattern_variant(pattern, variant)

    def _clear_slides(self):
        """Remove all existing slides from the unpacked template."""
        slides_dir = self.work_dir / 'ppt' / 'slides'
        rels_dir = slides_dir / '_rels'

        for f in list(slides_dir.glob('slide*.xml')):
            f.unlink()
        for f in list(rels_dir.glob('slide*.xml.rels')):
            f.unlink()

        # Remove notes
        notes_dir = self.work_dir / 'ppt' / 'notesSlides'
        if notes_dir.exists():
            shutil.rmtree(notes_dir)

        # Clean presentation.xml
        pres_path = self.work_dir / 'ppt' / 'presentation.xml'
        pres = pres_path.read_text('utf-8')
        pres = re.sub(r'<p:sldIdLst>.*?</p:sldIdLst>',
                       '<p:sldIdLst>\n  </p:sldIdLst>', pres, flags=re.DOTALL)
        pres_path.write_text(pres, 'utf-8')

        # Clean presentation.xml.rels
        prels_path = self.work_dir / 'ppt' / '_rels' / 'presentation.xml.rels'
        prels = prels_path.read_text('utf-8')
        prels = re.sub(r'\s*<Relationship[^>]*Target="slides/[^"]*"[^>]*/>', '', prels)
        prels_path.write_text(prels, 'utf-8')

        # Clean [Content_Types].xml
        ct_path = self.work_dir / '[Content_Types].xml'
        ct = ct_path.read_text('utf-8')
        ct = re.sub(r'\s*<Override[^>]*PartName="/ppt/slides/[^"]*"[^>]*/>', '', ct)
        ct = re.sub(r'\s*<Override[^>]*PartName="/ppt/notesSlides/[^"]*"[^>]*/>', '', ct)
        ct_path.write_text(ct, 'utf-8')

    def _next_shape_id(self):
        self._shape_id_counter += 1
        return self._shape_id_counter

    def _add_slide_from_layout(self, layout_key):
        """Add a slide from a named layout. Returns the slide XML path."""
        layout_file = LAYOUTS.get(layout_key)
        if not layout_file:
            raise ValueError(f"Unknown layout: {layout_key}. Available: {list(LAYOUTS.keys())}")

        output = add_slide(str(self.work_dir), layout_file)

        # Parse the output for slide filename and sldId
        slide_match = re.search(r'Created (slide\d+\.xml)', output)
        rid_match = re.search(r'r:id="(rId\d+)"', output)

        if not slide_match or not rid_match:
            raise RuntimeError(f"Failed to add slide: {output}")

        slide_file = slide_match.group(1)
        rid = rid_match.group(1)
        self._slide_count += 1
        slide_id = 255 + self._slide_count

        # Add to presentation.xml sldIdLst
        pres_path = self.work_dir / 'ppt' / 'presentation.xml'
        pres = pres_path.read_text('utf-8')
        sld_id_xml = f'<p:sldId id="{slide_id}" r:id="{rid}"/>'
        pres = pres.replace('</p:sldIdLst>', f'    {sld_id_xml}\n  </p:sldIdLst>')
        pres_path.write_text(pres, 'utf-8')

        slide_path = self.work_dir / 'ppt' / 'slides' / slide_file
        self._slide_layouts[str(slide_path)] = layout_key
        return slide_path

    def _extract_title_placeholder_bounds(self, layout_key):
        """Resolve the title placeholder bounds for a layout from template geometry."""
        theme = self._theme_config or _ACTIVE_THEME_CONFIG or {}
        layout_geom = theme.get('layout_geometry', {})
        info = layout_geom.get(str(layout_key).lower().replace('_', '-'), {})
        placeholders = info.get('placeholders', []) or []
        slide_w = float(theme.get('width', 13.33))
        slide_h = float(theme.get('height', 7.50))
        for ph in placeholders:
            if ph.get('type') not in ('title', 'ctrTitle'):
                continue
            bounds = ph.get('bounds', {})
            if not bounds:
                continue
            x_pct = float(bounds.get('x_pct', 0.0))
            y_pct = float(bounds.get('y_pct', 0.0))
            w_pct = float(bounds.get('w_pct', 0.0))
            h_pct = float(bounds.get('h_pct', 0.0))
            return {
                'x': round(slide_w * x_pct / 100.0, 2),
                'y': round(slide_h * y_pct / 100.0, 2),
                'w': round(slide_w * w_pct / 100.0, 2),
                'h': round(slide_h * h_pct / 100.0, 2),
            }
        return None

    def _build_slide_contract(self, slide, layout_key):
        """Build the runtime slide contract for a slide/layout pair."""
        behavior = semantic_behavior(layout_key, theme_config=self._theme_config)
        bounds = dict(content_bounds(layout_key))
        title_bounds = self._extract_title_placeholder_bounds(layout_key)
        return {
            'slide': str(slide),
            'layout_key': layout_key,
            'behavior': behavior,
            'bounds': bounds,
            'content_bg': _content_surface_bg(layout_key),
            'source_spec': _source_spec_for_bounds(bounds, behavior),
            'title_bounds': title_bounds,
        }

    def _remember_slide_contract(self, slide, layout_key):
        """Cache and activate the runtime slide contract for a slide."""
        contract = self._build_slide_contract(slide, layout_key)
        self._slide_contracts[str(slide)] = contract
        self._current_slide_bg = contract['content_bg']
        _CURRENT_SLIDE_CONTEXT['deck'] = self
        _CURRENT_SLIDE_CONTEXT['slide'] = slide
        return contract

    def current_slide_contract(self, slide=None):
        """Return the cached slide contract for a slide."""
        slide = slide or _CURRENT_SLIDE_CONTEXT.get('slide')
        if slide is None:
            raise RuntimeError("No active slide context. Add a slide first.")
        slide_key = str(slide)
        contract = self._slide_contracts.get(slide_key)
        if contract is not None:
            return contract
        layout_key = self._slide_layouts.get(slide_key)
        if not layout_key:
            raise KeyError(f"Slide layout not tracked for {slide}")
        return self._remember_slide_contract(slide, layout_key)

    def _validate_title(self, title, layout_key):
        """Fail early when a title breaks the layout contract."""
        import warnings as _warnings
        behavior = semantic_behavior(layout_key, theme_config=self._theme_config)
        raw = str(title or '')
        collapsed = re.sub(r'\s+', ' ', raw.replace('\v', ' ').replace('\n', ' ')).strip()
        if not collapsed:
            raise ValueError("Title cannot be empty.")
        if '\n' in raw and not behavior.get('allow_title_line_breaks'):
            role = behavior.get('semantic_role') or layout_key
            raise ValueError(f"Line breaks are not allowed in titles for {role}.")
        limit = int(behavior.get('title_max_chars') or 0)
        if limit and len(collapsed) > limit:
            role = behavior.get('semantic_role') or layout_key
            raise ValueError(
                f"Title is too long for {role} ({len(collapsed)} chars, max {limit}). "
                f"Shorten or choose a layout with a wider title treatment."
            )
        # Warn on near-limit titles so the agent can proactively shorten
        if limit and len(collapsed) > limit - 5:
            role = behavior.get('semantic_role') or layout_key
            _warnings.warn(
                f"Title near limit for {role} ({len(collapsed)}/{limit} chars). "
                f"Consider shortening to avoid 3-line wrap.",
                UserWarning, stacklevel=3,
            )
        return collapsed

    def _guard_authored_region(self, slide, x, y, w, h, kind='content'):
        """Prevent manual content from landing in disallowed zones."""
        contract = self.current_slide_contract(slide)
        behavior = contract['behavior']
        role = behavior.get('semantic_role') or contract['layout_key']
        bounds = contract['bounds']

        if behavior.get('image_only'):
            raise ValueError(
                f"{role} is a title+image layout. Use fill_picture()/add_stock_image() "
                f"and do not place manual {kind} content."
            )
        if behavior.get('statement') or behavior.get('structural') or not behavior.get('body_allowed', True):
            raise ValueError(
                f"{role} does not allow manual body content. Remove the manual {kind} "
                f"and rely on the layout or a supported pattern."
            )

        tol_x = 0.05
        tol_y = 0.05
        tol_bottom = 0.15
        min_x = bounds['x']
        max_x = bounds['x_end']
        if not behavior.get('protected_title_panel'):
            min_x = bounds.get('full_x', bounds['x'])
            max_x = bounds.get('full_x_end', bounds['x_end'])

        if x < min_x - tol_x:
            raise ValueError(
                f"{role} reserves the title/accent panel. Start {kind} at x>={bounds['x']:.2f}\"."
            )
        if y < bounds['y'] - tol_y:
            raise ValueError(
                f"{kind.capitalize()} starts at y={y:.2f}\" above the safe content band "
                f"for {role} (starts at {bounds['y']:.2f}\")."
            )
        if x + w > max_x + tol_x:
            raise ValueError(
                f"{kind.capitalize()} extends past the safe content width for {role}. "
                f"Keep x+w <= {max_x:.2f}\"."
            )
        if y + h > bounds['y_end'] + tol_bottom:
            raise ValueError(
                f"{kind.capitalize()} extends below the safe content zone for {role}. "
                f"Keep y+h <= {bounds['y_end']:.2f}\"."
            )

    def slide_contract(self, slide=None):
        """Public method wrapper for the current slide contract."""
        return self.current_slide_contract(slide)

    def _get_spTree(self, slide_path):
        """Parse a slide and return its spTree element + tree."""
        tree = ET.parse(slide_path)
        root = tree.getroot()
        spTree = root.find('.//p:cSld/p:spTree', NS)
        return tree, spTree

    def _add_placeholder(self, slide_path, ph_type, text, ph_idx=None):
        """Add a placeholder shape to a slide."""
        tree, spTree = self._get_spTree(slide_path)
        idx_attr = f' idx="{ph_idx}"' if ph_idx else ''
        escaped = _escape_xml(text)

        # Handle multi-line text using line breaks within one paragraph
        lines = escaped.split('\n')
        if len(lines) == 1:
            paras = f'<a:p><a:r><a:rPr lang="en-US" dirty="0"/><a:t>{lines[0]}</a:t></a:r></a:p>'
        else:
            runs = []
            for j, line in enumerate(lines):
                runs.append(f'<a:r><a:rPr lang="en-US" dirty="0"/><a:t>{line}</a:t></a:r>')
                if j < len(lines) - 1:
                    runs.append('<a:br/>')
            paras = '<a:p>' + ''.join(runs) + '</a:p>'

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="ph_{ph_type}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            f'<p:nvPr><p:ph type="{ph_type}"{idx_attr}/></p:nvPr>'
            f'</p:nvSpPr>'
            f'<p:spPr/>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'{paras}'
            f'</p:txBody></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_textbox(self, slide_path, text, x, y, w, h, sz=12, color=None, bold=False, align='left', valign='top', autofit=False, vertical=False):
        """Add a free text box to a slide.

        Args:
            autofit: If True, PowerPoint auto-shrinks text to fit the box
                     (prevents wrapping/overflow). Use for labels in tight spaces.
            vertical: If True, rotates text 270 deg (reads bottom-to-top).
                      Useful for row labels on process maps and matrices.
        """
        color = _sanitize_color(color)
        tree, spTree = self._get_spTree(slide_path)
        text_xml = _make_text_runs(text, sz=sz, color=color, bold=bold, align=align)
        anchor = _valign_map(valign)
        autofit_xml = '<a:normAutofit/>' if autofit else ''
        vert_attr = ' vert="vert270"' if vertical else ''

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="TB{self._shape_id_counter}"/>'
            f'<p:cNvSpPr txBox="1"/>'
            f'<p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:noFill/>'
            f'</p:spPr>'
            f'<p:txBody>'
            f'<a:bodyPr wrap="square" rtlCol="0" anchor="{anchor}"{vert_attr}>{autofit_xml}</a:bodyPr>'
            f'<a:lstStyle/>'
            f'{text_xml}'
            f'</p:txBody></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_bullet_box(self, slide_path, items, x, y, w, h, sz=12, color=None, spc_after=0):
        color = color or COLORS['DARK_TEXT']
        """Add a bulleted text box to a slide."""
        color = _sanitize_color(color)
        tree, spTree = self._get_spTree(slide_path)
        bullets_xml = _make_bullet_runs(items, sz=sz, color=color, spc_after=spc_after)
        anchor = "t"  # bullets always top-aligned

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="BL{self._shape_id_counter}"/>'
            f'<p:cNvSpPr txBox="1"/>'
            f'<p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:noFill/>'
            f'</p:spPr>'
            f'<p:txBody>'
            f'<a:bodyPr wrap="square" rtlCol="0" anchor="{anchor}"/>'
            f'<a:lstStyle/>'
            f'{bullets_xml}'
            f'</p:txBody></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_rectangle(self, slide_path, x, y, w, h, fill_color, line_color=None, line_width=None, shadow=False, dash=None):
        """Add a filled rectangle shape.

        Args:
            dash: Border dash style. Options: 'dash', 'dot', 'dashDot', 'lgDash'.
                  Requires line_color to be set.
        """
        fill_color = _sanitize_color(fill_color)
        line_color = _sanitize_color(line_color)
        tree, spTree = self._get_spTree(slide_path)

        line_xml = '<a:ln><a:noFill/></a:ln>'
        if line_color:
            lw = int((line_width or 1) * 12700)
            dash_xml = f'<a:prstDash val="{dash}"/>' if dash else ''
            line_xml = f'<a:ln w="{lw}">{dash_xml}<a:solidFill><a:srgbClr val="{line_color}"/></a:solidFill></a:ln>'

        effect_xml = _shadow_xml() if shadow else ''

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="R{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'
            f'{line_xml}'
            f'{effect_xml}'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_oval(self, slide_path, x, y, w, h, fill_color,
                  line_color=None, line_width=None):
        """Add a filled oval/circle shape with optional border."""
        fill_color = _sanitize_color(fill_color)
        tree, spTree = self._get_spTree(slide_path)

        if line_color:
            line_color = _sanitize_color(line_color)
            lw = int((line_width or 1) * 12700)
            ln_xml = f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{line_color}"/></a:solidFill></a:ln>'
        else:
            ln_xml = '<a:ln><a:noFill/></a:ln>'

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="O{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'
            f'{ln_xml}'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_line(self, slide_path, x, y, w, h, color=None, width=2, dash=None):
        """Add a line shape.

        Args:
            dash: Line dash style. None for solid. Options: 'dash', 'dot',
                  'dashDot', 'lgDash', 'lgDashDot', 'sysDash', 'sysDot'
        """
        color = color or _smart_accent()
        color = _sanitize_color(color)
        tree, spTree = self._get_spTree(slide_path)
        lw = int(width * 12700)
        dash_xml = f'<a:prstDash val="{dash}"/>' if dash else ''

        # OOXML requires non-negative cx/cy. For right-to-left or bottom-to-top
        # lines, use flipH/flipV and adjust origin to top-left of bounding box.
        flip_h = w < 0
        flip_v = h < 0
        if flip_h:
            x = x + w
            w = abs(w)
        if flip_v:
            y = y + h
            h = abs(h)
        flip_parts = []
        if flip_h:
            flip_parts.append('flipH="1"')
        if flip_v:
            flip_parts.append('flipV="1"')
        flip_attr = (' ' + ' '.join(flip_parts)) if flip_parts else ''

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="L{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm{flip_attr}><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
            f'<a:ln w="{lw}">{dash_xml}<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln>'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _stage_image_media(self, image_ref):
        """Write an image path or bundled stock-image ref into ppt/media/."""
        image = _resolve_image_source(image_ref)
        media_dir = self.work_dir / 'ppt' / 'media'
        media_dir.mkdir(parents=True, exist_ok=True)
        ext = image['ext']
        img_hash = hashlib.md5(image['bytes']).hexdigest()[:8]
        media_name = f'img_{self._next_shape_id()}_{img_hash}{ext}'
        dest = media_dir / media_name
        dest.write_bytes(image['bytes'])

        # Ensure content type is registered.
        content_type = image['content_type']
        ct_path = self.work_dir / '[Content_Types].xml'
        ct = ct_path.read_text('utf-8')
        ext_no_dot = ext.lstrip('.')
        if f'Extension="{ext_no_dot}"' not in ct:
            ct = ct.replace('</Types>',
                f'<Default Extension="{ext_no_dot}" ContentType="{content_type}"/></Types>')
            ct_path.write_text(ct, 'utf-8')
        return media_name

    def _add_image(self, slide_path, image_path, x, y, w, h):
        """Add an image to a slide.

        Copies the image into ppt/media/, adds a relationship in the slide's
        .rels file, and inserts a <p:pic> element into the slide spTree.
        """
        media_name = self._stage_image_media(image_path)

        # Add relationship in slide .rels
        slide_name = Path(slide_path).name
        rels_dir = self.work_dir / 'ppt' / 'slides' / '_rels'
        rels_path = rels_dir / f'{slide_name}.rels'
        rels = rels_path.read_text('utf-8')

        # Find next rId
        rids = [int(m) for m in re.findall(r'Id="rId(\d+)"', rels)]
        next_rid = max(rids) + 1 if rids else 2
        rid = f'rId{next_rid}'

        new_rel = (
            f'<Relationship Id="{rid}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
            f'Target="../media/{media_name}"/>'
        )
        rels = rels.replace('</Relationships>', f'{new_rel}</Relationships>')
        rels_path.write_text(rels, 'utf-8')

        # Add <p:pic> element to slide
        tree, spTree = self._get_spTree(slide_path)
        sid = self._next_shape_id()

        pic_xml = (
            f'<p:pic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvPicPr>'
            f'<p:cNvPr id="{sid}" name="Img{sid}"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
            f'<p:nvPr/>'
            f'</p:nvPicPr>'
            f'<p:blipFill>'
            f'<a:blip r:embed="{rid}"/>'
            f'<a:stretch><a:fillRect/></a:stretch>'
            f'</p:blipFill>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'</p:spPr></p:pic>'
        )

        spTree.append(ET.fromstring(pic_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_gradient_rectangle(self, slide_path, x, y, w, h,
                                color1, color2, angle=0):
        """Add a rectangle with linear gradient fill.

        Args:
            color1: Start color hex (e.g. '29BA74')
            color2: End color hex (e.g. '3EAD92')
            angle: Gradient angle in degrees (0=left-to-right, 90=top-to-bottom)
        """
        color1 = _sanitize_color(color1)
        color2 = _sanitize_color(color2)
        tree, spTree = self._get_spTree(slide_path)
        # OOXML gradient angle is in 60,000ths of a degree
        ang = int(angle * 60000)

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="GR{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:gradFill>'
            f'<a:gsLst>'
            f'<a:gs pos="0"><a:srgbClr val="{color1}"/></a:gs>'
            f'<a:gs pos="100000"><a:srgbClr val="{color2}"/></a:gs>'
            f'</a:gsLst>'
            f'<a:lin ang="{ang}" scaled="1"/>'
            f'</a:gradFill>'
            f'<a:ln><a:noFill/></a:ln>'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_shape(self, slide_path, preset, x, y, w, h, fill_color,
                   line_color=None, line_width=None, rotation=None,
                   adjustments=None, theme_color=None):
        """Add an arbitrary preset geometry shape.

        Args:
            preset: OOXML preset name (e.g. 'chevron', 'rightArrow',
                    'roundRect', 'triangle', 'pentagon', 'hexagon',
                    'diamond', 'star5', 'blockArc')
            rotation: Rotation in degrees (clockwise)
            adjustments: Dict of OOXML adjustment guides, e.g.
                        {'adj1': 10800000, 'adj2': 0, 'adj3': 25000}
                        Values are raw OOXML units (60000 per degree for angles,
                        1/50000 for ratios).
            theme_color: Scheme color name (e.g. 'accent1', 'tx2') — used
                        instead of fill_color when set. Produces <a:schemeClr>
                        instead of <a:srgbClr>.
        """
        tree, spTree = self._get_spTree(slide_path)

        rot_attr = ''
        if rotation:
            rot_attr = f' rot="{int(rotation * 60000)}"'

        # Fill: theme color or RGB
        if theme_color:
            fill_xml = f'<a:solidFill><a:schemeClr val="{theme_color}"/></a:solidFill>'
        else:
            fill_color = _sanitize_color(fill_color)
            fill_xml = f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'

        line_xml = '<a:ln><a:noFill/></a:ln>'
        if line_color:
            line_color = _sanitize_color(line_color)
            lw = int((line_width or 1) * 12700)
            line_xml = f'<a:ln w="{lw}"><a:solidFill><a:srgbClr val="{line_color}"/></a:solidFill></a:ln>'

        # Adjustment values
        if adjustments:
            av_items = ''.join(f'<a:gd name="{k}" fmla="val {v}"/>' for k, v in adjustments.items())
            av_xml = f'<a:avLst>{av_items}</a:avLst>'
        else:
            av_xml = '<a:avLst/>'

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="S{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm{rot_attr}><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="{preset}">{av_xml}</a:prstGeom>'
            f'{fill_xml}'
            f'{line_xml}'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def _add_rounded_rectangle(self, slide_path, x, y, w, h, fill_color,
                                radius=16667, line_color=None, line_width=None, shadow=False, dash=None):
        """Add a rounded rectangle. radius is in 1/50000 units (default ~1/3 rounding).

        Args:
            dash: Border dash style. Options: 'dash', 'dot', 'dashDot', 'lgDash'.
                  Requires line_color to be set.
        """
        fill_color = _sanitize_color(fill_color)
        line_color = _sanitize_color(line_color)
        tree, spTree = self._get_spTree(slide_path)

        line_xml = '<a:ln><a:noFill/></a:ln>'
        if line_color:
            lw = int((line_width or 1) * 12700)
            dash_xml = f'<a:prstDash val="{dash}"/>' if dash else ''
            line_xml = f'<a:ln w="{lw}">{dash_xml}<a:solidFill><a:srgbClr val="{line_color}"/></a:solidFill></a:ln>'

        effect_xml = _shadow_xml() if shadow else ''

        sp_xml = (
            f'<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{self._next_shape_id()}" name="RR{self._shape_id_counter}"/>'
            f'<p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="roundRect"><a:avLst>'
            f'<a:gd name="adj" fmla="val {radius}"/>'
            f'</a:avLst></a:prstGeom>'
            f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'
            f'{line_xml}'
            f'{effect_xml}'
            f'</p:spPr></p:sp>'
        )

        spTree.append(ET.fromstring(sp_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    # ============================================================
    # PUBLIC API - STRUCTURAL SLIDES
    # ============================================================

    def add_title_slide(self, title, subtitle='', date='', detail=False):
        """Add a BCG title slide with real template background/logo.

        Args:
            title: Main presentation title
            subtitle: Subtitle or description
            date: Date string, e.g., '3 March 2026'
            detail: Use detail template variant (white BG, smaller titles)
        """
        key = 'd_title' if detail else 'title'
        self._validate_title(title, key)
        slide = self._add_slide_from_layout(key)
        self._add_placeholder(slide, 'ctrTitle', title)
        if subtitle:
            self._add_placeholder(slide, 'subTitle', subtitle, ph_idx='1')
        if date:
            self._add_placeholder(slide, 'body', date, ph_idx='12')
        self._remember_slide_contract(slide, key)
        return slide

    def add_section_divider(self, title, layout='section_header_box', detail=False):
        """Add a section divider slide.

        Args:
            title: Section title (white text on green)
            layout: Layout variant. Options: 'section_header_box' (default),
                    'section_header_line', 'big_statement_green'
            detail: Use detail template variant (white BG, smaller titles)
        """
        key = _DETAIL_MAP.get(layout, layout) if detail else layout
        self._validate_title(title, key)
        slide = self._add_slide_from_layout(key)
        self._add_placeholder(slide, 'title', title)
        self._remember_slide_contract(slide, key)
        return slide

    def add_disclaimer(self, detail=False):
        """Add the BCG disclaimer slide (pre-filled from template layout)."""
        key = 'd_disclaimer' if detail else 'disclaimer'
        slide = self._add_slide_from_layout(key)
        self._remember_slide_contract(slide, key)
        return slide

    def add_end_slide(self, detail=False):
        """Add the BCG end slide with logo (from template layout)."""
        key = 'd_end' if detail else 'end'
        slide = self._add_slide_from_layout(key)
        self._remember_slide_contract(slide, key)
        return slide

    def add_quote_slide(self, quote, attribution='', detail=False):
        """Add a quote / big statement slide."""
        key = 'd_big_statement_green' if detail else 'big_statement_green'
        self._validate_title(quote, key)
        slide = self._add_slide_from_layout(key)
        self._add_placeholder(slide, 'title', quote)
        if attribution:
            self._add_textbox(slide, f'-- {attribution}', 0.69, 6.2, 11.96, 0.4,
                            sz=14, color=COLORS['MED_GRAY'])
        self._remember_slide_contract(slide, key)
        return slide

    def add_agenda_slide(self, items, detail=False):
        """Add an agenda slide with numbered items.

        Args:
            items: List of agenda item strings
            detail: Use detail template variant (white BG, smaller titles)
        """
        slide = self._add_slide_from_layout('d_title_only' if detail else 'title_only')
        self._add_placeholder(slide, 'title', 'Agenda')
        # Both templates use DARK_TEXT -- detail (white BG) and standard (gray BG)
        text_color = COLORS['DARK_TEXT']
        line_color = COLORS.get('OUTLINE', 'E0E0E0')
        y = CONTENT_START_Y + 0.1
        for i, item in enumerate(items, 1):
            # Green number
            self._add_textbox(slide, str(i), 0.69, y, 0.4, 0.4,
                            sz=16, color=COLORS['BCG_GREEN'], bold=True)
            # Item text
            self._add_textbox(slide, item, 1.2, y, 11.0, 0.4,
                            sz=14, color=text_color)
            # Separator line
            self._add_line(slide, 0.69, y + 0.45, 11.96, 0, color=line_color, width=0.5)
            y += 0.6
        self._remember_slide_contract(slide, 'd_title_only' if detail else 'title_only')
        return slide

    # ============================================================
    # PUBLIC API - CONTENT SLIDES
    # ============================================================

    def add_content_slide(self, title, source='Source: BCG analysis', layout='title_only', detail=True):
        """Add a content slide with action title. Returns slide path for
        further content additions using the shape methods.

        Args:
            title: Action title (complete sentence stating the 'so what')
            source: Source line text
            layout: Base layout key (default 'title_only')
            detail: Use detail template variant (default True — white BG,
                    smaller title, more content space). Pass False for
                    key-message slides that need the bold gray layout.
        """
        key = _DETAIL_MAP.get(layout, layout) if detail else layout
        self._validate_title(title, key)
        slide = self._add_slide_from_layout(key)
        contract = self._remember_slide_contract(slide, key)

        # Track the slide's background color for contrast-aware defaults
        self._add_placeholder(slide, 'title', title)
        if source and contract.get('source_spec'):
            # Source line sits at y≈6.74" — always on the slide master
            # background, never on an accent panel. Use SLIDE_BG for
            # contrast computation even on split layouts where
            # _current_slide_bg reflects the accent panel fill.
            source_bg = COLORS.get('SLIDE_BG', 'F2F2F2')
            bg_rgb = _hex_to_rgb(source_bg) or (242, 242, 242)
            bg_lum = _relative_luminance(bg_rgb)
            if bg_lum < 0.4:
                source_color = COLORS.get('WHITE', 'FFFFFF')
            else:
                # Use MED_GRAY but ensure WCAG 3:1 contrast against slide bg
                candidate = COLORS['MED_GRAY']
                ratio = _contrast_ratio(candidate, source_bg)
                if ratio < 3.0:
                    # Fall back to DARK_TEXT which always passes on light backgrounds
                    source_color = COLORS.get('DARK_TEXT', '575757')
                else:
                    source_color = candidate
            source_spec = contract['source_spec']
            self._add_textbox(slide, source,
                            source_spec['x'], source_spec['y'],
                            source_spec['w'], source_spec['h'],
                            sz=10, color=source_color)
        return slide

    # Expose shape methods for content building
    def add_textbox(self, slide, text, x, y, w, h, sz=12, color=None, bold=False, align='left', valign='top', autofit=False, vertical=False, bg=None):
        """Add a text box to a slide.

        Args:
            autofit: If True, PowerPoint auto-shrinks text to fit the box.
            vertical: If True, rotates text 270 deg.
            bg: Background color this text sits on. Text color auto-adjusts
                for contrast (dark bg -> white text, light bg -> dark text).
        """
        self._guard_authored_region(slide, x, y, w, h, kind='text')
        if color is None:
            # Use explicit bg if provided, else use the current slide's background
            effective_bg = bg or getattr(self, '_current_slide_bg', None)
            color = _auto_text_color(effective_bg)
        self._add_textbox(slide, text, x, y, w, h, sz=sz, color=color, bold=bold, align=align, valign=valign, autofit=autofit, vertical=vertical)

    def add_bullets(self, slide, items, x, y, w, h, sz=12, color=None, spc_after=6, bg=None):
        """Add a bulleted list to a slide.

        Args:
            spc_after: Space after each bullet in points (6=default, 0=tight, 10=spacious).
            bg: Background color this text sits on. Auto-adjusts text color.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='bullets')
        if color is None:
            effective_bg = bg or getattr(self, '_current_slide_bg', None)
            color = _auto_text_color(effective_bg)
        self._add_bullet_box(slide, items, x, y, w, h, sz=sz, color=color, spc_after=spc_after)

    def add_rectangle(self, slide, x, y, w, h, fill_color=None, line_color=None, line_width=None, shadow=False, dash=None):
        """Add a filled rectangle.

        Args:
            fill_color: Fill hex. If None, auto-selects based on slide background
                        (LIGHT_BG on white slides, WHITE on gray slides).

        Auto-adds a border if the fill is too similar to the slide background
        (prevents invisible shapes).
        """
        self._guard_authored_region(slide, x, y, w, h, kind='shape')
        if fill_color is None:
            slide_bg = getattr(self, '_current_slide_bg', COLORS.get('SLIDE_BG', 'F2F2F2'))
            bg_lum = _relative_luminance(_hex_to_rgb(slide_bg) or (242, 242, 242))
            fill_color = COLORS.get('LIGHT_BG', 'F2F2F2') if bg_lum > 0.85 else COLORS.get('WHITE', 'FFFFFF')
        # Auto-border: if shape fill matches slide bg, add subtle border
        if line_color is None and w > 0.5 and h > 0.3:
            slide_bg = getattr(self, '_current_slide_bg', COLORS.get('SLIDE_BG', 'F2F2F2'))
            if slide_bg and fill_color:
                ratio = _contrast_ratio(fill_color, slide_bg)
                if ratio < 1.3:  # Nearly identical to slide bg
                    line_color = COLORS.get('OUTLINE', 'E0E0E0')
                    line_width = line_width or 0.75
        self._add_rectangle(slide, x, y, w, h, fill_color, line_color, line_width, shadow=shadow, dash=dash)

    def add_oval(self, slide, x, y, w, h, fill_color=None,
                 line_color=None, line_width=None, theme_color=None):
        """Add a filled oval/circle with optional border.

        Args:
            fill_color: Fill hex color (or None if using theme_color).
            line_color: Border color (hex). None = no border.
            line_width: Border width in points (default 1).
            theme_color: Scheme color name (e.g. 'accent1') — overrides fill_color.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='shape')
        if theme_color:
            # Use the shape method which supports theme colors
            self._add_shape(slide, 'ellipse', x, y, w, h, fill_color,
                           line_color, line_width, theme_color=theme_color)
        else:
            self._add_oval(slide, x, y, w, h, fill_color, line_color, line_width)

    def add_line(self, slide, x, y, w, h=0, color=None, width=2, dash=None):
        """Add a line.

        Args:
            dash: Line style. None=solid, 'dash', 'dot', 'dashDot',
                  'lgDash', 'lgDashDot', 'sysDash', 'sysDot'
        """
        color = color or _smart_accent()
        self._add_line(slide, x, y, w, h, color, width, dash)

    def add_number_badge(self, slide, number, x, y, size=0.4, fill_color=None):
        """Add a colored circle with number/letter, auto-contrast text."""
        self._guard_authored_region(slide, x, y, size, size, kind='badge')
        fill_color = fill_color or _smart_accent()
        txt_color = _auto_text_color(fill_color)
        self._add_oval(slide, x, y, size, size, fill_color)
        self._add_textbox(slide, str(number), x, y, size, size,
                         sz=int(size * 30), color=txt_color, bold=True,
                         align='center', valign='middle')

    def add_image(self, slide, image_path, x, y, w, h):
        """Add an image (PNG, JPG, EMF, etc.) to a slide.

        Args:
            slide: Slide path from add_content_slide() etc.
            image_path: Filesystem path or stock-image ref from select_image_by_tags()
            x, y, w, h: Position and size in inches
        """
        self._guard_authored_region(slide, x, y, w, h, kind='image')
        self._add_image(slide, image_path, x, y, w, h)

    def fill_picture(self, slide, image_path):
        """Fill the layout's picture placeholder with an image.

        The slide must use a layout that has a picture placeholder:
        green_half, green_two_third, d_green_half, or d_green_two_third.

        Unlike add_image() which places images freely, this method creates a
        proper placeholder-bound picture that PowerPoint recognizes as filling
        the layout's designated image zone.

        Args:
            slide: Slide path from _add_slide_from_layout()
            image_path: Filesystem path or stock-image ref from select_image_by_tags()
        """
        slide_key = str(slide)
        layout = self._slide_layouts.get(slide_key)
        if not layout:
            raise ValueError(
                "Slide layout not tracked. Use a layout with a picture placeholder "
                "(green_half, green_two_third, d_green_half, d_green_two_third)."
            )

        ph_info = _PICTURE_PLACEHOLDERS.get(layout)
        if not ph_info:
            raise ValueError(
                f"Layout '{layout}' has no picture placeholder. "
                f"Use one of: {list(_PICTURE_PLACEHOLDERS.keys())}"
            )

        media_name = self._stage_image_media(image_path)

        # Add relationship in slide .rels
        slide_name = Path(slide).name
        rels_dir = self.work_dir / 'ppt' / 'slides' / '_rels'
        rels_path = rels_dir / f'{slide_name}.rels'
        rels = rels_path.read_text('utf-8')
        rids = [int(m) for m in re.findall(r'Id="rId(\d+)"', rels)]
        next_rid = max(rids) + 1 if rids else 2
        rid = f'rId{next_rid}'
        new_rel = (
            f'<Relationship Id="{rid}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
            f'Target="../media/{media_name}"/>'
        )
        rels = rels.replace('</Relationships>', f'{new_rel}</Relationships>')
        rels_path.write_text(rels, 'utf-8')

        # Calculate center-crop offsets (CSS background-size:cover behaviour).
        # srcRect l/t/r/b values are in 1/1000th of a percent (100000 = 100%).
        # We clip the source image so its aspect ratio matches the placeholder,
        # then stretch the cropped portion to fill -- no distortion.
        src_rect_attr = ''
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(image_path) as _img:
                img_w, img_h = _img.size
            ph_w, ph_h = ph_info['w'], ph_info['h']
            src_ratio = img_w / img_h
            dst_ratio = ph_w / ph_h
            if src_ratio > dst_ratio + 0.001:
                # Image wider than placeholder: crop left and right equally
                each = int((1 - dst_ratio / src_ratio) / 2 * 100000)
                src_rect_attr = f'l="{each}" t="0" r="{each}" b="0"'
            elif src_ratio < dst_ratio - 0.001:
                # Image taller than placeholder: crop top and bottom equally
                each = int((1 - src_ratio / dst_ratio) / 2 * 100000)
                src_rect_attr = f'l="0" t="{each}" r="0" b="{each}"'
        except Exception:
            pass  # PIL unavailable or unreadable -- fall back to plain stretch

        src_rect_xml = f'<a:srcRect {src_rect_attr}/>' if src_rect_attr else ''

        # Build <p:pic> with placeholder reference so PowerPoint fills the layout zone
        tree, spTree = self._get_spTree(slide)
        sid = self._next_shape_id()
        idx = ph_info['idx']
        x, y, w, h = ph_info['x'], ph_info['y'], ph_info['w'], ph_info['h']

        pic_xml = (
            f'<p:pic xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvPicPr>'
            f'<p:cNvPr id="{sid}" name="Picture {sid}"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
            f'<p:nvPr><p:ph type="pic" idx="{idx}"/></p:nvPr>'
            f'</p:nvPicPr>'
            f'<p:blipFill>'
            f'<a:blip r:embed="{rid}"/>'
            f'{src_rect_xml}'
            f'<a:stretch><a:fillRect/></a:stretch>'
            f'</p:blipFill>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'</p:spPr></p:pic>'
        )

        spTree.append(ET.fromstring(pic_xml))
        tree.write(slide, xml_declaration=True, encoding='UTF-8')

    def add_stock_image(self, slide, tags, match_mode='any'):
        """Select a stock image by tags and fill the slide's picture placeholder.

        Combines tag-based image lookup with fill_picture(). The slide must
        use a layout with a picture placeholder (green_half, green_two_third).

        Args:
            slide: Slide path from _add_slide_from_layout()
            tags: List of tag strings (e.g. ['logistics', 'warehouse'])
            match_mode: 'any' (at least one tag) or 'all' (every tag must match)

        Returns:
            True if an image was found and placed, False otherwise.
        """
        path = _select_image_from_tags(tags, match_mode)
        if path:
            self.fill_picture(slide, str(path))
            return True
        return False

    @staticmethod
    def select_image_by_tags(tags, match_mode='any'):
        """Look up a stock image by tags without placing it.

        Use this to check availability before choosing a layout, or to
        get the path for manual placement with add_image().

        Args:
            tags: List of tag strings (e.g. ['logistics', 'warehouse'])
            match_mode: 'any' or 'all'

        Returns:
            Asset ref string accepted by add_image()/fill_picture(), or None.
        """
        path = _select_image_from_tags(tags, match_mode)
        return str(path) if path else None

    def _ensure_icon_index(self):
        """Load and index the icon library once per deck (cached).

        Uses icons_bundle.json (single file, all icons zlib-compressed).
        """
        if hasattr(self, '_icon_idx'):
            return self._icon_idx
        icons_dir = find_skill_root() / 'assets' / 'icons'
        bundle_path = icons_dir / 'icons_bundle.json'
        if not bundle_path.exists():
            raise FileNotFoundError(f"Icon bundle not found at {bundle_path}")
        with open(bundle_path, 'r') as f:
            self._icon_bundle = json.load(f)
        raw = {'icons': {k: {'cx': v['cx'], 'cy': v['cy']} for k, v in self._icon_bundle['icons'].items()}}

        # Build token -> set of keys index
        token_map = {}
        key_tokens = {}
        display_map = {}
        for key, meta in raw['icons'].items():
            name = key.replace('bcgIcons_', '').replace('svg_', '')
            name_tokens = set(_tokenize(name))
            key_tokens[key] = name_tokens
            display_map[_normalize_icon_lookup(_display_icon_name(key, icon_type='icon'))] = key
            for tok in name_tokens:
                token_map.setdefault(tok, set()).add(key)
        self._icon_idx = {
            'raw': raw, 'dir': icons_dir,
            'token_map': token_map, 'key_tokens': key_tokens, 'display_map': display_map,
        }
        return self._icon_idx

    def _load_icon_xml(self, icon_key):
        """Load icon XML by key from icons_bundle.json (base64 + zlib)."""
        if icon_key not in self._icon_bundle.get('icons', {}):
            raise FileNotFoundError(f"Icon '{icon_key}' not found in bundle")
        return zlib.decompress(base64.b64decode(self._icon_bundle['icons'][icon_key]['z'])).decode('utf-8')

    def _rank_icon_candidates(self, query, idx, icon_type='icon'):
        """Return ranked icon candidates for a concept query.

        Scoring favors exact/prefix token hits, multi-token overlap, and
        full-name matches. Query tokens are expanded with a small synonym map
        for common consulting/business concepts.
        """
        query_tokens = _tokenize(query)
        query_joined = query.lower().replace(' ', '').replace('_', '')

        if not query_tokens and query_joined:
            query_tokens = [query_joined]
        query_tokens = _expand_icon_query_tokens(query_tokens)

        candidates = {}  # key -> score

        for qt in query_tokens:
            # Exact token matches (score x10)
            for key in idx['token_map'].get(qt, set()):
                candidates[key] = candidates.get(key, 0) + 10

            # Prefix matches (score x6)
            for tok, keys in idx['token_map'].items():
                if tok.startswith(qt) and tok != qt:
                    for key in keys:
                        candidates[key] = candidates.get(key, 0) + 6

        query_token_set = set(query_tokens)

        # Exact full-name match + token-overlap bonus
        for key, key_tokens in idx['key_tokens'].items():
            short = _display_icon_name(key, icon_type=icon_type).lower().replace('_', '')
            if short == query_joined:
                candidates[key] = candidates.get(key, 0) + 50
            elif query_joined in short:
                # Name substring bonus (score x3)
                candidates[key] = candidates.get(key, 0) + 3

            overlap = len(query_token_set & key_tokens)
            if overlap:
                candidates[key] = candidates.get(key, 0) + (overlap * 12)

        ranked = sorted(
            candidates.items(),
            key=lambda item: (-item[1], _display_icon_name(item[0], icon_type=icon_type).lower())
        )
        return ranked

    def search_icons(self, query, icon_type='icon', limit=8):
        """Return ranked icon candidates without placing an icon.

        This is the preferred discovery path when the exact icon name is
        unclear. It avoids ad hoc shell searches of the bundled icon files.

        Returns:
            list of dicts with keys: name, key, score
        """
        if icon_type == 'bug':
            idx = self._ensure_bug_index()
        else:
            idx = self._ensure_icon_index()

        ranked = self._rank_icon_candidates(query, idx, icon_type=icon_type)
        return [
            {
                'name': _display_icon_name(key, icon_type=icon_type),
                'key': key,
                'score': score,
            }
            for key, score in ranked[:max(limit, 1)]
        ]

    def _resolve_icon_key(self, query, idx, icon_type='icon'):
        """Resolve a query to a raw bundled icon key."""
        query = str(query or "").strip()
        if not query:
            raise ValueError("Icon query cannot be empty.")
        if query in idx['raw']['icons']:
            return query

        exact = idx.get('display_map', {}).get(_normalize_icon_lookup(query))
        if exact:
            return exact

        ranked = self._rank_icon_candidates(query, idx, icon_type=icon_type)
        if not ranked:
            available = sorted(
                _display_icon_name(k, icon_type=icon_type)
                for k in idx['raw']['icons']
            )[:30]
            label = "Bug icon" if icon_type == 'bug' else "Icon"
            raise ValueError(
                f"{label} '{query}' not found. Try: {', '.join(available[:15])}..."
            )
        return ranked[0][0]

    def resolve_icon(self, query, icon_type='icon'):
        """Resolve a concept query to the raw key used in the bundled icon set."""
        idx = self._ensure_bug_index() if icon_type == 'bug' else self._ensure_icon_index()
        return self._resolve_icon_key(query, idx, icon_type=icon_type)

    def _search_icon(self, query):
        """Token-based icon search with weighted scoring.

        Returns the best-matching icon key, or raises ValueError.
        """
        idx = self._ensure_icon_index()
        return self._resolve_icon_key(query, idx, icon_type='icon')

    def _ensure_bug_index(self):
        """Load and index the bug icon library once per deck (cached).
        Loads from bugs_bundle.json (single file, all icons zlib-compressed).
        """
        if hasattr(self, '_bug_idx'):
            return self._bug_idx
        
        # load the bundle of bug icons (base64 + zlib) from assets/bugs/bugs_bundle.json
        bugs_dir = find_skill_root() / 'assets' / 'bugs'
        bundle_path = bugs_dir / 'bugs_bundle.json'
        if not bundle_path.exists():
            raise FileNotFoundError(f"Bug bundle not found at {bundle_path}")
        
        # Load the bundle and extract metadata (cx, cy) for each icon.  
        with open(bundle_path, 'r') as f:
            self._bug_bundle = json.load(f)
        raw = {'icons': {k: {'cx': v['cx'], 'cy': v['cy']} for k, v in self._bug_bundle['icons'].items()}}

        # Build token -> set of keys index
        token_map = {}
        key_tokens = {}
        display_map = {}
        for key in raw['icons']:
            # Strip trailing _bug suffix for cleaner token matching
            name = key[:-4] if key.endswith('_bug') else key
            name_tokens = set(_tokenize(name))
            key_tokens[key] = name_tokens
            display_map[_normalize_icon_lookup(_display_icon_name(key, icon_type='bug'))] = key
            for tok in name_tokens:
                token_map.setdefault(tok, set()).add(key)
        self._bug_idx = {
            'raw': raw, 'dir': bugs_dir,
            'token_map': token_map, 'key_tokens': key_tokens, 'display_map': display_map,
        }
        return self._bug_idx

    def _load_bug_xml(self, bug_key):
        """Load bug XML by key from bugs_bundle.json (base64 + zlib)."""
        if bug_key not in self._bug_bundle.get('icons', {}):
            raise FileNotFoundError(f"Bug icon '{bug_key}' not found in bundle")
        return zlib.decompress(base64.b64decode(self._bug_bundle['icons'][bug_key]['z'])).decode('utf-8')
        

    def _search_bug(self, query):
        """Token-based bug icon search. Same scoring as _search_icon."""
        idx = self._ensure_bug_index()
        return self._resolve_icon_key(query, idx, icon_type='bug')

    def resolve_icon_type(self, line_count):
        """Resolve whether to use 'icon' or 'bug' based on slide text density.

        Count all text lines on the slide (title + bullets + table rows +
        body paragraphs + chart annotations) and pass the total here.
        Returns 'icon' for sparse slides, 'bug' for dense slides.

        Rule:
            line_count < 15  → 'icon'  (full BCG vector icon in oval circle)
            line_count >= 15 → 'bug'   (compact variant from assets/bugs/)
            When in doubt    → 'icon'

        Args:
            line_count: Total text lines on the slide (int).

        Returns:
            'icon' or 'bug'
        """
        if line_count >= 15:
            return 'bug'
        return 'icon'

    def add_icon(self, slide, icon_name, x, y, size=0.75, color=None, icon_type='icon'):
        """Add a BCG vector icon to a slide.

        1,069 OOXML vector icons with token-based search. Icons render as
        crisp vector graphics at any size with optional color override.

        Args:
            slide: Slide path from add_content_slide(), etc.
            icon_name: Icon name or concept (e.g. 'Target', 'Strategy',
                       'Global', 'Heart', 'Chess King'). Case-insensitive.
                       Uses token matching with stop-word filtering.
            x, y: Position in inches.
            size: Display size in inches (icons are scaled uniformly).
            color: Override fill color hex (e.g. '29BA74'). If None,
                   uses the icon's original theme colors.
            icon_type: 'icon' (default) or 'bug'. Resolved via
                       deck.resolve_icon_type(line_count) before calling.
                       'bug' searches assets/bugs/bugs_bundle.json; 'icon'
                       uses assets/icons/icons_bundle.json.

        Search by name or concept using the bundled icon libraries.
        """
        self._guard_authored_region(slide, x, y, size, size, kind='icon')
        if icon_type == 'bug':
            icon_key = self.resolve_icon(icon_name, icon_type='bug')
            icon_xml = self._load_bug_xml(icon_key)
            orig = self._ensure_bug_index()['raw']['icons'][icon_key]
        else:
            icon_key = self.resolve_icon(icon_name, icon_type='icon')
            icon_xml = self._load_icon_xml(icon_key)
            orig = self._ensure_icon_index()['raw']['icons'][icon_key]
            
        orig_cx = orig['cx']
        orig_cy = orig['cy']

        # Calculate scaled dimensions (uniform scaling based on size)
        aspect = orig_cx / orig_cy if orig_cy else 1
        if aspect >= 1:
            w_emu = int(size * 914400)
            h_emu = int(w_emu / aspect)
        else:
            h_emu = int(size * 914400)
            w_emu = int(h_emu * aspect)

        # Update the group shape transform
        x_emu = emu(x)
        y_emu = emu(y)

        # Parse and modify the icon XML
        icon_tree = ET.fromstring(icon_xml)

        # Update outer group transform
        grp_xfrm = icon_tree.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
        if grp_xfrm is not None:
            off = grp_xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}off')
            ext = grp_xfrm.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
            if off is not None:
                off.set('x', x_emu)
                off.set('y', y_emu)
            if ext is not None:
                ext.set('cx', str(w_emu))
                ext.set('cy', str(h_emu))

        # Override color if specified
        color = _sanitize_color(color)
        if color:
            for fill in icon_tree.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill'):
                # Replace scheme colors with explicit RGB
                scheme = fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                if scheme is not None:
                    fill.remove(scheme)
                    srgb = ET.SubElement(fill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    srgb.set('val', color)
                # Replace existing srgbClr
                srgb = fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    srgb.set('val', color)

        # Update shape IDs to avoid conflicts
        for elem in icon_tree.iter():
            cNvPr_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr'
            if elem.tag == cNvPr_tag:
                elem.set('id', str(self._next_shape_id()))

        # Add to slide
        tree, spTree = self._get_spTree(slide)
        spTree.append(icon_tree)
        tree.write(slide, xml_declaration=True, encoding='UTF-8')

    def add_gradient_rectangle(self, slide, x, y, w, h, color1, color2, angle=0):
        """Add a rectangle with linear gradient fill.

        Args:
            slide: Slide path
            color1: Start color hex (e.g. '29BA74')
            color2: End color hex (e.g. '3EAD92')
            angle: Gradient angle in degrees (0=left-to-right, 90=top-to-bottom)
        """
        self._guard_authored_region(slide, x, y, w, h, kind='shape')
        self._add_gradient_rectangle(slide, x, y, w, h, color1, color2, angle)

    def add_shape(self, slide, preset, x, y, w, h, fill_color=None,
                  line_color=None, line_width=None, rotation=None,
                  adjustments=None, theme_color=None):
        """Add a preset geometry shape.

        Args:
            preset: Shape name - 'chevron', 'rightArrow', 'triangle',
                    'pentagon', 'hexagon', 'diamond', 'star5', 'heart',
                    'cloud', 'gear6', 'gear9', 'blockArc', etc.
            rotation: Rotation in degrees (clockwise)
            adjustments: Dict of OOXML adjustment guides for the preset.
            theme_color: Scheme color name (e.g. 'accent1') — overrides fill_color.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='shape')
        self._add_shape(slide, preset, x, y, w, h, fill_color,
                       line_color, line_width, rotation, adjustments, theme_color)

    def add_arc_sector(self, slide, x, y, w, h, start_deg, sweep_deg,
                       thickness_pct=25, fill_color=None, theme_color=None,
                       line_color='FFFFFF', line_width=0.5):
        """Add a donut arc sector (blockArc preset shape).

        Args:
            x, y, w, h: Bounding box in inches (the arc is inscribed in this box)
            start_deg: Start angle in degrees (0=right, 90=bottom, -90=top)
            sweep_deg: Sweep angle in degrees (positive = clockwise)
            thickness_pct: Arc thickness as percentage of radius (0-50). Default 25.
            fill_color: Fill hex color (or None if using theme_color)
            theme_color: Scheme color name (e.g. 'accent1') — overrides fill_color.
            line_color: Border color hex. Default white.
            line_width: Border width in points. Default 0.5.
        """
        # blockArc adjustments:
        #   adj1 = start angle in 60000ths of a degree
        #   adj2 = end angle in 60000ths of a degree
        #   adj3 = thickness as 1/50000 of radius (50000 = full radius, 0 = no thickness)
        end_deg = start_deg + sweep_deg
        adj1 = int(start_deg * 60000) % 21600000
        adj2 = int(end_deg * 60000) % 21600000
        adj3 = int(thickness_pct * 1000)  # pct -> 1/50000 units (25% = 25000)

        self._add_shape(
            slide, 'blockArc', x, y, w, h, fill_color,
            line_color=line_color, line_width=line_width,
            adjustments={'adj1': adj1, 'adj2': adj2, 'adj3': adj3},
            theme_color=theme_color,
        )

    def add_rounded_rectangle(self, slide, x, y, w, h, fill_color=None,
                               radius=16667, line_color=None, line_width=None, shadow=False, dash=None):
        """Add a rounded rectangle.

        Args:
            fill_color: Fill hex. If None, auto-selects based on slide background
                        (LIGHT_BG on white slides, WHITE on gray slides).
            radius: Corner rounding in 1/50000 units. 16667=subtle, 50000=pill shape.
            shadow: Add subtle drop shadow (for card backgrounds).
            dash: Border dash style ('dash', 'dot', 'dashDot', 'lgDash'). Requires line_color.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='shape')
        if fill_color is None:
            slide_bg = getattr(self, '_current_slide_bg', COLORS.get('SLIDE_BG', 'F2F2F2'))
            bg_lum = _relative_luminance(_hex_to_rgb(slide_bg) or (242, 242, 242))
            fill_color = COLORS.get('LIGHT_BG', 'F2F2F2') if bg_lum > 0.85 else COLORS.get('WHITE', 'FFFFFF')
        self._add_rounded_rectangle(slide, x, y, w, h, fill_color,
                                     radius, line_color, line_width, shadow=shadow, dash=dash)

    def add_card(self, slide, x, y, w, h, fill_color=None, shadow=True,
                 accent_color=None, accent_position='top', radius=5000,
                 line_color=None, line_width=None, dash=None):
        """Add a card: rounded rectangle with optional accent bar and shadow.

        Cards are the primary container for content blocks in BCG slides.
        Shadow defaults to True (unlike plain rectangles).

        Args:
            fill_color: Card background color. If None, auto-selects based on
                        the slide background: LIGHT_BG (F2F2F2) on white/detail
                        slides, WHITE on gray/standard slides.
            shadow: Add subtle drop shadow (default: True).
            accent_color: Color for thin accent bar (e.g., COLORS['BCG_GREEN']).
                          None = no accent bar.
            accent_position: 'top' (0.06" bar at top) or 'left' (0.08" bar at left).
            radius: Corner rounding (default 5000 = subtle).
            line_color: Border color. Required for dash to take effect.
            line_width: Border width in points (default 1).
            dash: Border dash style ('dash', 'dot', 'dashDot', 'lgDash'). Requires line_color.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='card')
        if fill_color is None:
            # Auto-select card fill based on slide background for contrast
            slide_bg = getattr(self, '_current_slide_bg', COLORS.get('SLIDE_BG', 'F2F2F2'))
            bg_lum = _relative_luminance(_hex_to_rgb(slide_bg) or (242, 242, 242))
            if bg_lum > 0.85:
                # White/light background (detail slides) → light gray cards
                fill_color = COLORS.get('LIGHT_BG', 'F2F2F2')
            else:
                # Gray/darker background (standard slides) → white cards
                fill_color = COLORS.get('WHITE', 'FFFFFF')
        self._add_rounded_rectangle(slide, x, y, w, h, fill_color,
                                     radius=radius, line_color=line_color,
                                     line_width=line_width, shadow=shadow, dash=dash)
        if accent_color:
            if accent_position == 'top':
                self._add_rectangle(slide, x, y, w, 0.06, accent_color)
            elif accent_position == 'left':
                self._add_rectangle(slide, x, y, 0.08, h, accent_color)

    def add_label(self, slide, text, x, y, w, h=0.35, fill_color=None,
                  text_color=None, sz=11):
        """Add a colored label badge with centered text.

        Text color auto-detects from fill for best contrast.

        Args:
            text: Label text (kept short -- 1-3 words).
            fill_color: Background color (default: DARK_GREEN).
            text_color: Text color. Auto-detected from fill if not set.
            sz: Font size in points.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='label')
        fill_color = fill_color or COLORS['DARK_GREEN']
        text_color = text_color or _auto_text_color(fill_color)
        self._add_rectangle(slide, x, y, w, h, fill_color)
        self._add_textbox(slide, text, x, y, w, h,
                         sz=sz, color=text_color, bold=True,
                         align='center', valign='middle')

    def add_stamp(self, slide, text, x, y, w=1.8, h=0.35, fill_color=None,
                  text_color=None, sz=11):
        """Add a callout stamp badge (e.g., 'See Appendix', 'Draft').

        Text color auto-detects from fill for best contrast.

        Args:
            text: Stamp text (1-3 words).
            fill_color: Background color (default: STAMP_RED).
            text_color: Text color. Auto-detected from fill if not set.
            sz: Font size in points.
        """
        self._guard_authored_region(slide, x, y, w, h, kind='stamp')
        fill_color = fill_color or COLORS['STAMP_RED']
        text_color = text_color or _auto_text_color(fill_color)
        self._add_rectangle(slide, x, y, w, h, fill_color)
        self._add_textbox(slide, text, x, y, w, h,
                         sz=sz, color=text_color, bold=True,
                         align='center', valign='middle')

    def add_section_breadcrumb(self, slide, text, x=0.69, y=1.60, w=2.5, h=0.40,
                               fill_color=None, text_color=None, sz=10):
        """Add a parallelogram-shaped section breadcrumb below the title.

        Useful for showing which section of the deck the current slide belongs to.
        Positioned just above the content area by default.

        Args:
            text: Section name (e.g., 'Market Analysis').
            fill_color: Background color (default: DARK_GREEN).
            text_color: Text color (default: white).
            sz: Font size in points.
        """
        fill_color = fill_color or COLORS['DARK_GREEN']
        text_color = text_color or _auto_text_color(fill_color)
        self._add_shape(slide, 'parallelogram', x, y, w, h, fill_color=fill_color)
        self._add_textbox(slide, text, x, y, w, h,
                         sz=sz, color=text_color, bold=True,
                         align='center', valign='middle')

    def add_connector_dots(self, slide, x_start, x_end, y, n_dots=3,
                           dot_size=0.12, color=None):
        """Add evenly-spaced connector dots between two x positions.

        Useful for connecting process phases or timeline stages.
        """
        color = color or _smart_accent()
        spacing = (x_end - x_start) / (n_dots + 1)
        for i in range(1, n_dots + 1):
            cx = x_start + spacing * i - dot_size / 2
            cy = y - dot_size / 2
            self._add_oval(slide, cx, cy, dot_size, dot_size, color)

    def add_chevron_flow(self, slide, labels, x, y, w, h, colors=None):
        """Add a horizontal chevron flow (connected arrow shapes).

        Args:
            labels: List of strings for each chevron
            x, y: Start position in inches
            w: Total width for all chevrons
            h: Height of each chevron
            colors: List of fill colors (defaults to green gradient)
        """
        n = len(labels)
        if colors is None:
            # Build gradient from COLORS (theme-dynamic — works for BCG default and client ee4p)
            shades = [COLORS.get('DEEP_GREEN', COLORS['BCG_GREEN']),
                      COLORS['DARK_GREEN'],
                      COLORS['BCG_GREEN'],
                      COLORS.get('NAVY', COLORS.get('ACCENT5', COLORS['BCG_GREEN'])),
                      COLORS.get('DEEP_GREEN', COLORS['BCG_GREEN'])]
            colors = [shades[i % len(shades)] for i in range(n)]

        overlap = 0.15  # chevrons overlap slightly for connected look
        chev_w = (w + overlap * (n - 1)) / n
        for i in range(n):
            cx = x + i * (chev_w - overlap)
            self._add_shape(slide, 'chevron', cx, y, chev_w, h, colors[i])
            txt_color = _auto_text_color(colors[i])
            self._add_textbox(slide, labels[i], cx + 0.25, y, chev_w - 0.45, h,
                            sz=14, color=txt_color, bold=True,
                            align='center', valign='middle')

    # ============================================================
    # TABLES
    # ============================================================

    def _add_table(self, slide_path, data, x, y, w, h,
                   header=True, col_widths=None, col_align=None,
                   sz=11, header_color=None, stripe=True):
        """Add a native PowerPoint table to a slide.

        Args:
            data: List of rows, each row a list of cell strings.
                  First row is treated as header if header=True.
            x, y, w, h: Position and size in inches.
            header: If True, first row gets green fill + white bold text.
            col_widths: Proportional column widths, e.g. [2, 1, 1, 1].
                        If None, columns share width equally.
            col_align: Alignment per column, e.g. ['left', 'center', 'right'].
                       If None, first column is 'left', rest 'center'.
            sz: Body font size in points. Header uses sz+1.
            header_color: Header row fill color (hex). Default BCG Green.
            stripe: If True, alternating row fills (white / F2F2F2).
        """
        # Smart header color: use theme accent when available
        if header_color is None:
            header_color = _sanitize_color(_smart_accent())
        else:
            header_color = _sanitize_color(header_color)
        tree, spTree = self._get_spTree(slide_path)

        n_rows = len(data)
        n_cols = max(len(row) for row in data) if data else 0
        if n_rows == 0 or n_cols == 0:
            return

        w_emu = int(w * 914400)
        # Row height: use 0.40" per row as the target, capped to available height.
        # This prevents rows from stretching to fill the entire content zone
        # when the table has fewer rows than the zone can hold.
        target_row_h = 0.40  # inches
        max_row_h = h / n_rows
        actual_row_h = min(target_row_h, max_row_h)
        h_emu = int(actual_row_h * n_rows * 914400)
        row_h_emu = int(actual_row_h * 914400)

        # Column widths
        if col_widths:
            total = sum(col_widths)
            grid_cols = [int((cw / total) * w_emu) for cw in col_widths]
        else:
            grid_cols = [w_emu // n_cols] * n_cols

        # Column alignments
        if col_align is None:
            col_align = ['l'] + ['ctr'] * (n_cols - 1)
        else:
            _map = {'left': 'l', 'center': 'ctr', 'right': 'r', 'l': 'l', 'ctr': 'ctr', 'r': 'r'}
            col_align = [_map.get(a, 'l') for a in col_align]

        # Padding in EMU (0.1" left/right, 0.05" top/bottom)
        mar_lr = '91440'
        mar_tb = '45720'

        # Build grid XML
        grid_xml = ''.join(f'<a:gridCol w="{gw}"/>' for gw in grid_cols)

        # Build rows
        rows_xml = []
        for ri, row in enumerate(data):
            is_hdr = header and ri == 0
            cells_xml = []

            for ci in range(n_cols):
                text = str(row[ci]) if ci < len(row) else ''
                escaped = _escape_xml(text)

                # Text properties
                bold_attr = ' b="1"' if is_hdr else ''
                # Compute text color from the cell's background fill for contrast
                if is_hdr:
                    txt_color = _auto_text_color(header_color)
                elif stripe and ri % 2 == 0:
                    txt_color = text_on(COLORS.get('SURFACE_ALT', COLORS['LIGHT_BG']))
                else:
                    txt_color = text_on(COLORS['WHITE'])
                font_sz = pt_emu(sz + 1) if is_hdr else pt_emu(sz)
                align = col_align[ci] if ci < len(col_align) else 'l'

                # Cell fill
                if is_hdr:
                    fill = f'<a:solidFill><a:srgbClr val="{header_color}"/></a:solidFill>'
                elif stripe and ri % 2 == 0:
                    # Use SURFACE_ALT for stripe (theme-aware, not hardcoded F2F2F2)
                    stripe_clr = COLORS.get('SURFACE_ALT', COLORS['LIGHT_BG'])
                    fill = f'<a:solidFill><a:srgbClr val="{stripe_clr}"/></a:solidFill>'
                else:
                    # White fill for non-stripe rows (visible separation from slide bg)
                    fill = f'<a:solidFill><a:srgbClr val="{COLORS["WHITE"]}"/></a:solidFill>'

                # Borders: no borders (clean modern look) — row fills provide structure
                bdr = (
                    '<a:lnL w="0"><a:noFill/></a:lnL>'
                    '<a:lnR w="0"><a:noFill/></a:lnR>'
                    '<a:lnT w="0"><a:noFill/></a:lnT>'
                    '<a:lnB w="0"><a:noFill/></a:lnB>'
                )

                # Handle multi-line cell text
                lines = escaped.split('\n')
                paras = []
                for line in lines:
                    paras.append(
                        f'<a:p><a:pPr algn="{align}"/>'
                        f'<a:r><a:rPr lang="en-US" sz="{font_sz}" dirty="0"{bold_attr}>'
                        f'<a:solidFill><a:srgbClr val="{txt_color}"/></a:solidFill>'
                        f'<a:latin typeface="{FONT}"/><a:cs typeface="{FONT}"/>'
                        f'</a:rPr><a:t>{line}</a:t></a:r></a:p>'
                    )

                cells_xml.append(
                    f'<a:tc>'
                    f'<a:txBody>'
                    f'<a:bodyPr anchor="ctr" lIns="{mar_lr}" tIns="{mar_tb}" '
                    f'rIns="{mar_lr}" bIns="{mar_tb}"/>'
                    f'<a:lstStyle/>'
                    f'{"".join(paras)}'
                    f'</a:txBody>'
                    f'<a:tcPr marL="{mar_lr}" marR="{mar_lr}" marT="{mar_tb}" marB="{mar_tb}">'
                    f'{bdr}{fill}'
                    f'</a:tcPr>'
                    f'</a:tc>'
                )

            rows_xml.append(
                f'<a:tr h="{row_h_emu}">{"".join(cells_xml)}</a:tr>'
            )

        # Build graphicFrame
        sid = self._next_shape_id()
        frame_xml = (
            f'<p:graphicFrame '
            f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            f'<p:nvGraphicFramePr>'
            f'<p:cNvPr id="{sid}" name="Table {sid}"/>'
            f'<p:cNvGraphicFramePr><a:graphicFrameLocks noGrp="1"/></p:cNvGraphicFramePr>'
            f'<p:nvPr/>'
            f'</p:nvGraphicFramePr>'
            f'<p:xfrm>'
            f'<a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/>'
            f'</p:xfrm>'
            f'<a:graphic>'
            f'<a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">'
            f'<a:tbl>'
            f'<a:tblPr firstRow="1" bandRow="0"><a:tblStyle val="{{00000000-0000-0000-0000-000000000000}}"/><a:noFill/></a:tblPr>'
            f'<a:tblGrid>{grid_xml}</a:tblGrid>'
            f'{"".join(rows_xml)}'
            f'</a:tbl>'
            f'</a:graphicData>'
            f'</a:graphic>'
            f'</p:graphicFrame>'
        )

        spTree.append(ET.fromstring(frame_xml))
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    def add_table(self, slide, data, x=0.69, y=None, w=11.96, h=None,
                  header=True, col_widths=None, col_align=None,
                  sz=11, header_color=None, stripe=True):
        """Add a native PowerPoint table to a slide.

        Args:
            slide: Slide path from add_content_slide(), etc.
            data: List of rows. Each row is a list of cell values (strings).
                  First row is treated as header if header=True.
                  Example: [['Metric', 'Q1', 'Q2'], ['Revenue', '$2.1B', '$2.3B']]
            x, y, w: Position and width in inches. y defaults to CONTENT_START_Y.
            h: Total table height. If None, auto-calculated (0.4" per row).
            header: First row gets BCG Green fill + white bold text (default True).
            col_widths: Proportional widths, e.g. [3, 1, 1, 1]. None = equal.
            col_align: Alignment per column: 'left', 'center', 'right'.
                       None = first column left, rest center.
            sz: Body font size (header gets sz+1). Default 11.
            header_color: Header fill hex color. Default COLORS['BCG_GREEN'].
            stripe: Alternating row colors white/F2F2F2. Default True.
        """
        if y is None:
            y = self.current_slide_contract(slide)['bounds']['y']
        if h is None:
            h = len(data) * 0.4
        self._guard_authored_region(slide, x, y, w, h, kind='table')
        # header_color=None triggers smart default in _add_table
        self._add_table(slide, data, x, y, w, h,
                       header=header, col_widths=col_widths, col_align=col_align,
                       sz=sz, header_color=header_color, stripe=stripe)

    # ============================================================
    # SPEAKER NOTES
    # ============================================================

    def set_speaker_notes(self, slide, text):
        """Add speaker notes to a slide.

        Creates a notesSlide XML file linked to the given slide.
        Speaker notes appear in Presenter View and Notes Page view.

        Args:
            slide: Slide path from add_content_slide(), etc.
            text: Speaker notes text. Newlines create separate paragraphs.
                  Keep notes concise: 3-5 bullet points per slide.
        """
        slide_name = Path(slide).name
        slide_num = re.search(r'slide(\d+)', slide_name).group(1)
        notes_name = f'notesSlide{slide_num}.xml'

        # Create notesSlides directory
        notes_dir = self.work_dir / 'ppt' / 'notesSlides'
        notes_dir.mkdir(parents=True, exist_ok=True)
        rels_dir = notes_dir / '_rels'
        rels_dir.mkdir(parents=True, exist_ok=True)

        # Build notes paragraphs
        paragraphs = text.split('\n')
        paras_xml = ''
        for para in paragraphs:
            escaped = _escape_xml(para)
            paras_xml += (
                f'<a:p><a:r><a:rPr lang="en-US" dirty="0">'
                f'<a:latin typeface="{FONT}"/></a:rPr>'
                f'<a:t>{escaped}</a:t></a:r></a:p>'
            )

        # Build notes slide XML
        notes_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:cSld><p:spTree>'
            '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
            '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
            '<p:sp><p:nvSpPr>'
            f'<p:cNvPr id="2" name="Slide Image Placeholder 1"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr>'
            '<p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr>'
            '<p:spPr/></p:sp>'
            '<p:sp><p:nvSpPr>'
            f'<p:cNvPr id="3" name="Notes Placeholder 2"/>'
            '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
            '<p:spPr/>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>{paras_xml}</p:txBody>'
            '</p:sp></p:spTree></p:cSld></p:notes>'
        )

        notes_path = notes_dir / notes_name
        notes_path.write_text(notes_xml, 'utf-8')

        # Create notes .rels file linking back to the slide
        notes_rels = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            f'<Relationship Id="rId1" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" '
            f'Target="../slides/{slide_name}"/>'
            '</Relationships>'
        )
        (rels_dir / f'{notes_name}.rels').write_text(notes_rels, 'utf-8')

        # Add relationship from slide to notes
        slide_rels_dir = self.work_dir / 'ppt' / 'slides' / '_rels'
        slide_rels_path = slide_rels_dir / f'{slide_name}.rels'
        slide_rels = slide_rels_path.read_text('utf-8')

        # Find next rId
        rids = [int(m) for m in re.findall(r'Id="rId(\d+)"', slide_rels)]
        next_rid = max(rids) + 1 if rids else 2
        notes_rel = (
            f'<Relationship Id="rId{next_rid}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" '
            f'Target="../notesSlides/{notes_name}"/>'
        )
        slide_rels = slide_rels.replace('</Relationships>', f'{notes_rel}</Relationships>')
        slide_rels_path.write_text(slide_rels, 'utf-8')

        # Register notes content type
        ct_path = self.work_dir / '[Content_Types].xml'
        ct = ct_path.read_text('utf-8')
        if 'notesSlide+xml' not in ct:
            ct = ct.replace('</Types>',
                '<Override PartName="/ppt/notesSlides/' + notes_name + '" '
                'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                '</Types>')
        else:
            # Add this specific notes slide
            if f'/ppt/notesSlides/{notes_name}' not in ct:
                ct = ct.replace('</Types>',
                    '<Override PartName="/ppt/notesSlides/' + notes_name + '" '
                    'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                    '</Types>')
        ct_path.write_text(ct, 'utf-8')

    # ============================================================
    # NATIVE CHARTS (via python-pptx)
    # ============================================================

    def add_chart(self, slide, chart_type, categories, series,
                  x=0.69, y=None, w=11.96, h=4.0, **options):
        """Add a native editable chart to a slide.

        Charts are injected during save() via python-pptx. Requires:
            pip install python-pptx

        Args:
            slide: Slide path from add_content_slide(), etc.
            chart_type: One of:
                'bar'            - Clustered column chart
                'stacked_bar'    - Stacked column chart
                'bar_horizontal' - Clustered horizontal bars
                'line'           - Line chart with markers
                'pie'            - Pie chart
                'doughnut'       - Doughnut chart
                'area'           - Area chart
                'stacked_area'   - Stacked area chart
            categories: List of category labels ['Q1', 'Q2', 'Q3', 'Q4']
            series: List of dicts with 'name' and 'values':
                [{'name': 'Revenue', 'values': [10, 20, 30, 40]}, ...]
                Optional per-series keys:
                    point_colors  - Per-point hex fills for emphasis on single-series bars
                    line_dash     - Line style for line charts ('solid', 'dash', 'dashDot', 'lgDash')
            x, y, w, h: Position and size in inches.
                y defaults to CONTENT_START_Y (2.10")
            **options:
                colors: List of hex strings (e.g. ['197A56', '29BA74'])
                number_format: Value axis format (e.g. '$#,##0', '0%')
                data_labels: True to show data labels (default False)
                data_label_color: Hex color for data labels (e.g., 'FFFFFF' for white)
                data_label_position: 'outEnd', 'inEnd', 'ctr', 'inBase', 'bestFit', etc.
                legend: True/False to show legend (default True)
                gap_width: Bar gap percentage (default 150)
                hole_size: Doughnut hole percentage (default 50)
                smooth: True for smooth lines (default False)
                value_axis_min: Numeric minimum for value axis
                value_axis_max: Numeric maximum for value axis

        Returns:
            None (chart injected at save time)
        """
        if y is None:
            y = self.current_slide_contract(slide)['bounds']['y']
        self._guard_authored_region(slide, x, y, w, h, kind='chart')

        match = re.search(r'slide(\d+)\.xml', str(slide))
        if not match:
            raise ValueError(f"Cannot determine slide number from: {slide}")

        chart_series = []
        for s in series:
            series_spec = dict(s)
            series_spec['values'] = list(s['values'])
            if 'point_colors' in series_spec and series_spec['point_colors'] is not None:
                series_spec['point_colors'] = list(series_spec['point_colors'])
            chart_series.append(series_spec)

        self._chart_specs.append({
            'slide_index': int(match.group(1)) - 1,
            'chart_type': chart_type,
            'categories': list(categories),
            'series': chart_series,
            'x': x, 'y': y, 'w': w, 'h': h,
            'options': dict(options),
        })

    def _inject_charts(self, path):
        """Post-process saved PPTX to inject native editable charts."""
        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.util import Inches, Pt
            from pptx.enum.chart import (
                XL_CHART_TYPE, XL_LEGEND_POSITION,
                XL_LABEL_POSITION, XL_MARKER_STYLE,
            )
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError(
                "python-pptx is required for native charts. "
                "Install with: pip install python-pptx"
            )

        TYPE_MAP = {
            'bar':            XL_CHART_TYPE.COLUMN_CLUSTERED,
            'stacked_bar':    XL_CHART_TYPE.COLUMN_STACKED,
            'bar_horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
            'line':           XL_CHART_TYPE.LINE_MARKERS,
            'pie':            XL_CHART_TYPE.PIE,
            'doughnut':       XL_CHART_TYPE.DOUGHNUT,
            'area':           XL_CHART_TYPE.AREA,
            'stacked_area':   XL_CHART_TYPE.AREA_STACKED,
        }

        PALETTE = [
            RGBColor(0x19, 0x7A, 0x56),  # Dark Green
            RGBColor(0x29, 0xBA, 0x74),  # BCG Green
            RGBColor(0x84, 0xE3, 0x87),  # Light Green
            RGBColor(0xC2, 0xDD, 0x79),  # Lime
            RGBColor(0x68, 0x89, 0x10),  # Olive
            RGBColor(0x00, 0xBC, 0xD4),  # Cyan
            RGBColor(0x29, 0x5E, 0x7E),  # Navy
            RGBColor(0x3E, 0xAD, 0x92),  # Teal
        ]

        DARK_TEXT = RGBColor(0x57, 0x57, 0x57)
        AXIS_LINE = RGBColor(0xB1, 0xB1, 0xB7)

        LABEL_POSITION_MAP = {
            'outEnd': XL_LABEL_POSITION.OUTSIDE_END,
            'outsideEnd': XL_LABEL_POSITION.OUTSIDE_END,
            'inEnd': XL_LABEL_POSITION.INSIDE_END,
            'insideEnd': XL_LABEL_POSITION.INSIDE_END,
            'inBase': XL_LABEL_POSITION.INSIDE_BASE,
            'insideBase': XL_LABEL_POSITION.INSIDE_BASE,
            'ctr': XL_LABEL_POSITION.CENTER,
            'center': XL_LABEL_POSITION.CENTER,
            'bestFit': XL_LABEL_POSITION.BEST_FIT,
            'above': XL_LABEL_POSITION.ABOVE,
            'below': XL_LABEL_POSITION.BELOW,
            'left': XL_LABEL_POSITION.LEFT,
            'right': XL_LABEL_POSITION.RIGHT,
        }
        DASH_STYLE_MAP = {
            'solid': MSO_LINE_DASH_STYLE.SOLID,
            'dash': MSO_LINE_DASH_STYLE.DASH,
            'dashDot': MSO_LINE_DASH_STYLE.DASH_DOT,
            'dashDotDot': MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
            'lgDash': MSO_LINE_DASH_STYLE.LONG_DASH,
            'lgDashDot': MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
            'dot': MSO_LINE_DASH_STYLE.ROUND_DOT,
            'squareDot': MSO_LINE_DASH_STYLE.SQUARE_DOT,
        }

        def _rgb(value):
            if value is None:
                return None
            if isinstance(value, RGBColor):
                return value
            return RGBColor.from_string(str(value).lstrip('#'))

        prs = Presentation(path)

        for spec in self._chart_specs:
            slide = prs.slides[spec['slide_index']]
            opts = spec['options']
            ct = spec['chart_type']

            # Build chart data
            chart_data = CategoryChartData()
            chart_data.categories = spec['categories']
            for s in spec['series']:
                chart_data.add_series(s['name'], tuple(s['values']))

            xl_type = TYPE_MAP.get(ct, XL_CHART_TYPE.COLUMN_CLUSTERED)

            # Add chart shape to slide
            chart_frame = slide.shapes.add_chart(
                xl_type,
                Inches(spec['x']), Inches(spec['y']),
                Inches(spec['w']), Inches(spec['h']),
                chart_data,
            )
            chart = chart_frame.chart
            plot = chart.plots[0]

            # ---- No chart title (action title serves this role) ----
            chart.has_title = False

            # ---- Legend ----
            if opts.get('legend', True) and ct not in ('pie',):
                chart.has_legend = True
                leg = chart.legend
                leg.include_in_layout = False
                leg.position = XL_LEGEND_POSITION.BOTTOM
                leg.font.name = FONT
                leg.font.size = Pt(9)
                leg.font.color.rgb = DARK_TEXT
            else:
                chart.has_legend = False

            # ---- Colors ----
            user_colors = opts.get('colors')
            pal = []
            if user_colors:
                for c in user_colors:
                    pal.append(_rgb(c))
            else:
                pal = list(PALETTE)

            if ct in ('pie', 'doughnut'):
                # Color each data point
                s0 = plot.series[0]
                for i in range(len(spec['categories'])):
                    pt = s0.points[i]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = pal[i % len(pal)]
            else:
                # Color each series
                for i, ser in enumerate(plot.series):
                    color = pal[i % len(pal)]
                    point_colors = spec['series'][i].get('point_colors')
                    if ct in ('line',):
                        ser.format.line.color.rgb = color
                        ser.format.line.width = Pt(2)
                        ser.smooth = opts.get('smooth', False)
                        ser.marker.style = XL_MARKER_STYLE.CIRCLE
                        ser.marker.size = 6
                        ser.marker.format.fill.solid()
                        ser.marker.format.fill.fore_color.rgb = color
                        ser.marker.format.line.color.rgb = color
                        dash_style = spec['series'][i].get('line_dash')
                        if dash_style in DASH_STYLE_MAP:
                            ser.format.line.dash_style = DASH_STYLE_MAP[dash_style]
                    elif ct in ('area', 'stacked_area'):
                        ser.format.fill.solid()
                        ser.format.fill.fore_color.rgb = color
                    else:
                        ser.format.fill.solid()
                        ser.format.fill.fore_color.rgb = color

                        # Support per-point emphasis for single-series comparison charts.
                        if point_colors:
                            for j, point_color in enumerate(point_colors):
                                if point_color is None:
                                    continue
                                pt = ser.points[j]
                                pt.format.fill.solid()
                                pt.format.fill.fore_color.rgb = _rgb(point_color)

                # Convenience: allow colors=[...per-point...] on a single-series bar chart.
                if (
                    ct in ('bar', 'bar_horizontal')
                    and len(plot.series) == 1
                    and user_colors
                    and len(user_colors) == len(spec['categories'])
                ):
                    ser = plot.series[0]
                    for j, point_color in enumerate(user_colors):
                        pt = ser.points[j]
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = _rgb(point_color)

            # ---- Axis styling (non-pie charts) ----
            if ct not in ('pie', 'doughnut'):
                # Value axis
                va = chart.value_axis
                va.has_title = False
                va.has_major_gridlines = False
                va.has_minor_gridlines = False
                va.tick_labels.font.name = FONT
                va.tick_labels.font.size = Pt(9)
                va.tick_labels.font.color.rgb = DARK_TEXT
                # Hide value axis line
                try:
                    va.format.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    va.format.line.width = Pt(0)
                except Exception:
                    pass

                if 'number_format' in opts:
                    va.tick_labels.number_format = opts['number_format']
                    va.tick_labels.number_format_is_linked = False
                if 'value_axis_min' in opts:
                    va.minimum_scale = opts['value_axis_min']
                if 'value_axis_max' in opts:
                    va.maximum_scale = opts['value_axis_max']

                # Category axis
                ca = chart.category_axis
                ca.has_title = False
                ca.has_major_gridlines = False
                ca.tick_labels.font.name = FONT
                ca.tick_labels.font.size = Pt(9)
                ca.tick_labels.font.color.rgb = DARK_TEXT
                try:
                    ca.format.line.color.rgb = AXIS_LINE
                    ca.format.line.width = Pt(0.5)
                except Exception:
                    pass

            # ---- Bar gap width ----
            if ct in ('bar', 'stacked_bar', 'bar_horizontal'):
                try:
                    plot.gap_width = opts.get('gap_width', 150)
                except Exception:
                    pass

            # ---- Doughnut hole size ----
            if ct == 'doughnut':
                try:
                    # Access via XML since python-pptx may not expose directly
                    from lxml import etree
                    ns = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
                    hole_elem = plot._element.find(f'{{{ns}}}holeSize')
                    if hole_elem is not None:
                        hole_elem.set('val', str(opts.get('hole_size', 50)))
                except Exception:
                    pass

            # ---- Data labels ----
            if opts.get('data_labels', False):
                plot.has_data_labels = True
                dl = plot.data_labels
                dl.font.name = FONT
                dl.font.size = Pt(9)
                dl_color = opts.get('data_label_color', None)
                if dl_color:
                    dl.font.color.rgb = _rgb(dl_color)
                else:
                    dl.font.color.rgb = DARK_TEXT
                if 'number_format' in opts:
                    dl.number_format = opts['number_format']
                    dl.number_format_is_linked = False
                dl_position = opts.get('data_label_position')
                if dl_position in LABEL_POSITION_MAP:
                    dl.position = LABEL_POSITION_MAP[dl_position]

        prs.save(path)

    # ============================================================
    # SAVE
    # ============================================================

    @property
    def slide_count(self):
        """Return the number of slides added to this deck."""
        return self._slide_count

    def add_raw_xml_batch(self, slide, sp_xml_strings):
        """Append multiple raw OOXML shape elements in a single write.

        Use this for shapes that need custom geometry, scheme colors,
        or other OOXML features not supported by the preset shape methods.

        Args:
            slide: Slide path from add_content_slide()
            sp_xml_strings: List of complete <p:sp> XML strings with
                           proper namespace declarations.
        """
        tree, spTree = self._get_spTree(slide)
        for xml_str in sp_xml_strings:
            spTree.append(ET.fromstring(xml_str))
        tree.write(slide, xml_declaration=True, encoding='UTF-8')

    def save(self, output_path):
        """Clean, validate, and pack the presentation.

        Args:
            output_path: Output .pptx file path
        """
        if self.work_dir is None:
            raise RuntimeError("This deck was initialized for preflight only and cannot be saved.")

        # Clean
        clean(str(self.work_dir))

        # Pack
        pack(str(self.work_dir), output_path, original_path=self.template_path)

        # Inject native charts (requires python-pptx)
        if self._chart_specs:
            self._inject_charts(output_path)

        if os.path.exists(output_path):
            size = os.path.getsize(output_path)
            print(f"Saved: {output_path} ({size:,} bytes)")
        else:
            raise RuntimeError(f"Failed to save: {output_path}")

        return output_path

    def cleanup(self):
        """Remove the unpacked working directory when the deck is no longer needed."""
        if self.work_dir and Path(self.work_dir).exists():
            shutil.rmtree(self.work_dir, ignore_errors=True)
            self.work_dir = None


# ============================================================
# GRID HELPERS (same as bcg-base.js)
# ============================================================

def columns(n, total_width=11.96, start_x=0.69, gutter=0.31, gap=None):
    """Calculate equal-width column positions.

    Returns list of (x, w) tuples.

    Args:
        gap: Alias for gutter. If provided, overrides gutter.
    """
    if gap is not None:
        gutter = gap
    col_w = (total_width - (n - 1) * gutter) / n
    return [(start_x + i * (col_w + gutter), col_w) for i in range(n)]


def stack_y_positions(bounds, heights, gap=0.12, align='center', top_pad=0.0, bottom_pad=0.0):
    """Return top y positions for vertically stacked blocks inside bounds.

    This is the canonical helper for split-layout manual composition. It uses
    the active content band returned by `content_bounds()` rather than the full
    7.5" slide height.
    """
    heights = [float(h) for h in heights]
    if not heights:
        return []

    available_h = float(bounds['h']) - float(top_pad) - float(bottom_pad)
    required_h = sum(heights) + gap * max(len(heights) - 1, 0)
    if required_h > available_h + 1e-6:
        raise ValueError(
            f"Stack requires {required_h:.2f}\" of height, but only {available_h:.2f}\" is available "
            f"inside the safe content band."
        )

    if align == 'top':
        start_y = float(bounds['y']) + float(top_pad)
    elif align == 'bottom':
        start_y = float(bounds['y']) + float(top_pad) + (available_h - required_h)
    else:
        start_y = float(bounds['y']) + float(top_pad) + max(0.0, (available_h - required_h) / 2.0)

    positions = []
    cursor = start_y
    for height in heights:
        positions.append(round(cursor, 2))
        cursor += height + gap
    return positions
